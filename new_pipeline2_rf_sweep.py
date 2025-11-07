#!/usr/bin/env python3
# new_pipeline2_rf_sweep_with_cache_stream_incremental_fixed_with_cm_pptx_and_frfedback_multi_model.py
"""
Same as before but minimal edits so the script works for both BERT and other HF models
(e.g. "state-spaces/mamba-130m-hf").

Minimal changes:
 - Use AutoTokenizer / AutoModel and a generic MODEL_NAME variable
 - Replace references to `bert_model` with `base_model`
 - Added FINETUNE_ENABLED flag to enable/disable finetuning
Everything else is unchanged from the previous version, except that execution is placed
inside `main()` and wrapped by a try/except that emails success/failure via `special`.
"""
import os
import csv
import textwrap
import numpy as np
from tqdm import tqdm
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import confusion_matrix, f1_score
from transformers import AutoTokenizer, AutoModel, AutoConfig, AutoModelForCausalLM
import torch
import torch.nn as nn
import sys
import matplotlib.pyplot as plt
from datetime import datetime
from zoneinfo import ZoneInfo
# email & traceback imports for wrapper
import traceback
import special

# try to import python-pptx; if missing we'll raise a clear error when needed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

# ---------------- Config ----------------
device_list = [
    "irobot_roomba.json",
    "line_clova_wave.json",
    "nature_remo.json",
    "qrio_hub.json",
    "xiaomi_mijia_led.json",
    "powerelectric_wi-fi_plug.json",
    "planex_smacam_outdoor.json"
]
INPUT_FOLDER = "preprocessed_data_group_merged/ungrouped"

# Model selection: can be "bert-base-uncased" or "state-spaces/mamba-130m-hf" (or other HF model).
# MODEL_NAME = "bert-base-uncased"
# Example: MODEL_NAME = "state-spaces/mamba-130m-hf"
MODEL_NAME = "state-spaces/mamba-130m-hf"

# Detect Mamba-like model (minimal rule)
IS_MAMBA = ("mamba" in MODEL_NAME.lower()) or MODEL_NAME.startswith("state-spaces/mamba")

# Fraction of each device file to use (0.0 < DATA_USAGE_PCT <= 1.0)
# Set to 0.5 to use the first 50% of each device's data.
DATA_USAGE_PCT = 0.1  # default 1.0 (use everything); change to 0.5, 0.25, etc.


# ----- New flag: enable/disable finetuning entirely -----
# If True: perform finetuning (or reuse .ft if REUSE_FINETUNE_EMBS_IF_PRESENT allows).
# If False: skip finetuning and skip finetuned evaluation (only baseline results).
FINETUNE_ENABLED = True
# --------------------------------------------------------

# Sweep settings
START_PCT = 0.20   # final RF sweep will iterate RF train from 20% to 70%
END_PCT = 0.70
STEP_PCT = 0.02
RF_VAL_PCT = 0.30   # fixed: reserve last 30% per device for validation

# Fine-tune specification (legacy variable; reserved block is controlled by RESERVED_FINE_PCT)
FINE_PCT = 0.01     # first chunk per device used for fine-tuning (single student)

# Fixed reserved percent: always reserve the first 20% of each device (excluded from RF training)
RESERVED_FINE_PCT = 0.20

# If True: and valid finetuned embeddings (.ft.npy) exist for *all* devices, skip finetuning.
# If False: always (re)run finetuning and overwrite .ft embeddings.
REUSE_FINETUNE_EMBS_IF_PRESENT = True

# DEVICE (cuda if available)
DEVICE = torch.device('cuda' if torch.cuda.is_available() else 'cpu')

RF_N_ESTIMATORS = 500
RF_N_JOBS = -1

# Finetuning hyperparams
LR = 2e-5
EPOCHS = 3
FT_BATCH_SIZE = 16
LAMBDA_ANCHOR = 0.5
LAMBDA_PROTO = 0.5
LAMBDA_CLS = 0.5
LAMBDA_RF  = 0.5

# Embedding cache folder
EMB_DIR = "embeddings"

# Confusion matrix outputs
CM_DIR = os.path.join(EMB_DIR, "confusion_matrices")
CM_BASELINE_DIR = os.path.join(CM_DIR, "baseline")
CM_FINETUNED_DIR = os.path.join(CM_DIR, "finetuned")
os.makedirs(EMB_DIR, exist_ok=True)
os.makedirs(CM_DIR, exist_ok=True)
os.makedirs(CM_BASELINE_DIR, exist_ok=True)
os.makedirs(CM_FINETUNED_DIR, exist_ok=True)

# Misc / performance
BATCH_TOKEN_MAX_LEN = 512
EMBED_BATCH_SIZE = 32
PRINT_SAMPLE_LINES = True

# ---------------- Helpers ----------------
def short(s, width=120):
    return textwrap.shorten(s.replace("\n", "\\n"), width=width, placeholder="…")

def count_nonempty_lines(path):
    cnt = 0
    with open(path, 'r', encoding='utf-8') as f:
        for line in f:
            if line.strip():
                cnt += 1
    return cnt

def load_nonempty_lines(file_path):
    lines = []
    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            s = line.rstrip("\n")
            if s.strip():
                lines.append(s)
    return lines

# ---- New helper: unified batched embedding extraction (returns torch.Tensor on DEVICE) ----
def model_batch_embeddings_torch(batch_texts, tokenizer, model, max_length=BATCH_TOKEN_MAX_LEN, require_grad=False):
    """
    Returns torch.Tensor shape (B, H) on DEVICE.
    For Mamba: mean-pool last hidden state across non-padding tokens (attention mask).
    For BERT-like encoder models: return CLS token last_hidden_state[:,0,:].

    require_grad: if True, run the forward pass with gradients enabled (used during finetuning).
                  if False (default), run under no_grad() for faster inference / caching.
    """
    enc = tokenizer(batch_texts, return_tensors='pt', truncation=True, max_length=max_length, padding=True)
    enc = {k: v.to(DEVICE) for k, v in enc.items()}

    # enable grad only if requested (so we can use this helper both for inference caching and for training)
    with torch.set_grad_enabled(require_grad):
        out = model(**enc)

        if IS_MAMBA:
            # outputs.hidden_states[-1] shape: (B, seq_len, hidden)
            last_hidden = out.hidden_states[-1]  # (B, seq_len, hidden)
            mask = enc["attention_mask"].unsqueeze(-1)  # (B, seq_len, 1)
            summed = (last_hidden * mask).sum(dim=1)  # (B, hidden)
            counts = mask.sum(dim=1).clamp(min=1)  # (B, 1) avoid div by zero
            emb = summed / counts  # (B, hidden)
        else:
            # encoder-style CLS token
            emb = out.last_hidden_state[:, 0, :]  # (B, hidden)

    return emb  # torch tensor on DEVICE (may require_grad depending on require_grad)


def get_embedding_batch(sentences, tokenizer, model, max_length=512):
    """
    Compute embeddings for a list of sentences in one forward pass (batched).
    Returns numpy array shape (len(sentences), hidden_size).
    """
    emb_t = model_batch_embeddings_torch(sentences, tokenizer, model, max_length=max_length)
    return emb_t.cpu().numpy()

def interleave_round_robin(per_device_lists, devices_order, device_label_map):
    pointers = {d: 0 for d in devices_order}
    total_remaining = sum(len(per_device_lists[d]) for d in devices_order)
    combined = []
    labels = []
    while total_remaining > 0:
        for d in devices_order:
            idx = pointers[d]
            lst = per_device_lists[d]
            if idx < len(lst):
                combined.append(lst[idx])
                labels.append(device_label_map[d])
                pointers[d] += 1
                total_remaining -= 1
    return combined, labels

# ---------------- Prepare environment (model/tokenizer) ----------------
print("Loading tokenizer and model (this may take a moment)...")
tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, use_fast=True)
cfg = AutoConfig.from_pretrained(MODEL_NAME)
cfg.output_hidden_states = True  # ensure hidden states are returned

# Minimal branching for Mamba vs encoder models:
if IS_MAMBA:
    # follow the first script's pattern: causal-lm model class + float16 dtype
    base_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, torch_dtype=torch.float16).to(DEVICE)
else:
    base_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg).to(DEVICE)
base_model.eval()

device_label_map = {dev: idx for idx, dev in enumerate(device_list)}
device_short_names = [d.replace('.json', '') for d in device_list]

# ---------------- Precompute / load lines and set up caches ----------------
per_device_lines = {}
per_device_total_lines = {}

for device in device_list:
    file_name = device.replace('.json', '.json.txt')
    file_path = os.path.join(INPUT_FOLDER, file_name)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Missing file: {file_path}")

    total_lines = count_nonempty_lines(file_path)
    if total_lines == 0:
        raise ValueError(f"No non-empty lines in {file_path}")
    per_device_total_lines[device] = total_lines

    safe_name = device.replace('.json', '')
    lines_cache_path = os.path.join(EMB_DIR, f"{safe_name}.lines.txt")

    if os.path.exists(lines_cache_path):
        lines = load_nonempty_lines(lines_cache_path)
        if len(lines) != total_lines:
            print(f"  lines cache mismatch for {device}: cached={len(lines)} vs file={total_lines}. Recreating.")
            lines = load_nonempty_lines(file_path)
            with open(lines_cache_path, 'w', encoding='utf-8') as wf:
                for L in lines:
                    wf.write(L + "\n")
    else:
        lines = load_nonempty_lines(file_path)
        with open(lines_cache_path, 'w', encoding='utf-8') as wf:
            for L in lines:
                wf.write(L + "\n")

    per_device_lines[device] = lines

print("\nAll device line counts:")
sum_device_count=0
for d in device_list:
    print(f"  {d}: total non-empty lines = {per_device_total_lines[d]}")
    sum_device_count += per_device_total_lines[d]
print(f"Total lines across all devices: {sum_device_count}")
print("")

# ---------------- Streaming embedding computation & caching ----------------
from numpy.lib.format import open_memmap

def compute_and_cache_embeddings_stream(device, batch_size=EMBED_BATCH_SIZE):
    """
    Compute baseline embeddings (base_model) streaming and save .npy/.txt
    Uses unified model_batch_embeddings_torch helper (so Mamba vs BERT handled).
    """
    safe_name = device.replace('.json', '')
    lines = per_device_lines[device]
    n = len(lines)
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb.npy")
    emb_txt_path = os.path.join(EMB_DIR, f"{safe_name}.emb.txt")

    if os.path.exists(emb_npy_path):
        try:
            candidate = np.load(emb_npy_path, mmap_mode='r')
            if candidate.ndim == 2 and candidate.shape[0] == n and candidate.shape[1] == base_model.config.hidden_size:
                print(f"  Existing .npy for {device} looks valid; reusing: {emb_npy_path}")
                return candidate
            else:
                try:
                    os.remove(emb_npy_path)
                except OSError:
                    pass
                if os.path.exists(emb_txt_path):
                    try:
                        os.remove(emb_txt_path)
                    except OSError:
                        pass
        except Exception:
            try:
                os.remove(emb_npy_path)
            except OSError:
                pass
            if os.path.exists(emb_txt_path):
                try:
                    os.remove(emb_txt_path)
                except OSError:
                    pass

    hidden_size = base_model.config.hidden_size
    print(f"Streaming-computing embeddings for {device} ({n} lines). Writing .npy -> {emb_npy_path}")
    mm = open_memmap(emb_npy_path, dtype='float32', mode='w+', shape=(n, hidden_size))
    txt_f = open(emb_txt_path, 'w', encoding='utf-8')

    base_model.eval()
    with torch.no_grad():
        idx = 0
        for start in tqdm(range(0, n, batch_size), desc=f"Embedding (stream) {safe_name}"):
            batch_texts = lines[start:start+batch_size]
            batch_embs = get_embedding_batch(batch_texts, tokenizer, base_model, max_length=BATCH_TOKEN_MAX_LEN)
            B = batch_embs.shape[0]
            for b in range(B):
                row = batch_embs[b].astype(np.float32)
                mm[idx, :] = row
                txt_f.write(" ".join(f"{float(v):.6e}" for v in row) + "\n")
                idx += 1
    txt_f.flush()
    txt_f.close()
    mm.flush()
    arr = np.load(emb_npy_path, mmap_mode='r')
    print(f"  Finished {device}. .npy saved at {emb_npy_path} shape={arr.shape}")
    return arr

def compute_and_cache_embeddings_stream_for_model(device, model, suffix=".ft", batch_size=EMBED_BATCH_SIZE):
    """
    Compute embeddings with a given model (e.g., student) and save .npy/.txt with suffix.
    Uses unified helper to ensure Mamba vs BERT consistency.
    """
    safe_name = device.replace('.json', '')
    lines = per_device_lines[device]
    n = len(lines)
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.npy")
    emb_txt_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.txt")

    if os.path.exists(emb_npy_path):
        try:
            candidate = np.load(emb_npy_path, mmap_mode='r')
            if candidate.ndim == 2 and candidate.shape[0] == n and candidate.shape[1] == model.config.hidden_size:
                print(f"  Existing .npy (suffix={suffix}) for {device} looks valid; reusing: {emb_npy_path}")
                return candidate
            else:
                try:
                    os.remove(emb_npy_path)
                except OSError:
                    pass
                if os.path.exists(emb_txt_path):
                    try:
                        os.remove(emb_txt_path)
                    except OSError:
                        pass
        except Exception:
            try:
                os.remove(emb_npy_path)
            except OSError:
                pass
            if os.path.exists(emb_txt_path):
                try:
                    os.remove(emb_txt_path)
                except OSError:
                    pass

    hidden_size = model.config.hidden_size
    print(f"Streaming-computing embeddings (model suffix={suffix}) for {device} ({n} lines). Writing .npy -> {emb_npy_path}")
    mm = open_memmap(emb_npy_path, dtype='float32', mode='w+', shape=(n, hidden_size))
    txt_f = open(emb_txt_path, 'w', encoding='utf-8')

    model.eval()
    with torch.no_grad():
        idx = 0
        for start in tqdm(range(0, n, batch_size), desc=f"Embedding (stream){suffix} {safe_name}"):
            batch_texts = lines[start:start+batch_size]
            emb_t = model_batch_embeddings_torch(batch_texts, tokenizer, model, max_length=BATCH_TOKEN_MAX_LEN)
            batch_embs = emb_t.cpu().numpy()
            B = batch_embs.shape[0]
            for b in range(B):
                row = batch_embs[b].astype(np.float32)
                mm[idx, :] = row
                txt_f.write(" ".join(f"{float(v):.6e}" for v in row) + "\n")
                idx += 1
    txt_f.flush()
    txt_f.close()
    mm.flush()
    arr = np.load(emb_npy_path, mmap_mode='r')
    print(f"  Finished {device} (suffix={suffix}). .npy saved at {emb_npy_path} shape={arr.shape}")
    return arr

def load_embeddings_memmap_for_device(device, suffix=""):
    """
    Loads .npy file (if valid) or returns None so caller can compute it.
    """
    safe_name = device.replace('.json', '')
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.npy")
    total_lines = per_device_total_lines[device]
    if os.path.exists(emb_npy_path):
        try:
            arr = np.load(emb_npy_path, mmap_mode='r')
            if arr.ndim == 1:
                hidden_size = base_model.config.hidden_size
                if arr.size % hidden_size == 0:
                    rows = arr.size // hidden_size
                    arr = arr.reshape((rows, hidden_size))
            if arr.ndim == 2 and arr.shape[0] == total_lines:
                return arr
            else:
                return None
        except Exception:
            return None
    else:
        return None

# ---------------- Compute missing baseline embeddings up-front ----------------
def compute_missing_embeddings(devices):
    to_compute = []
    for device in devices:
        emb_arr = load_embeddings_memmap_for_device(device, suffix="")
        if emb_arr is None:
            to_compute.append(device)
        else:
            print(f"Skipping {device}: baseline .npy exists and looks valid ({emb_arr.shape[0]} rows).")
    for device in to_compute:
        compute_and_cache_embeddings_stream(device, batch_size=EMBED_BATCH_SIZE)

# ---------------- helper: check finetuned embeddings exist ----------------
def all_finetuned_embeddings_exist(devices, suffix=".ft"):
    hidden_size = base_model.config.hidden_size
    for device in devices:
        emb_npy_path = os.path.join(EMB_DIR, f"{device.replace('.json','')}.emb{suffix}.npy")
        if not os.path.exists(emb_npy_path):
            return False
        try:
            arr = np.load(emb_npy_path, mmap_mode='r')
            if arr.ndim == 1:
                if arr.size % hidden_size == 0:
                    rows = arr.size // hidden_size
                    arr = arr.reshape((rows, hidden_size))
                else:
                    return False
            if arr.ndim != 2:
                return False
            if arr.shape[0] != per_device_total_lines[device]:
                return False
        except Exception:
            return False
    return True

# ---------------- plotting / confusion / pptx helpers ----------------
def normalize_cm_rows(cm):
    cm = cm.astype(float)
    row_sums = cm.sum(axis=1, keepdims=True)
    with np.errstate(divide='ignore', invalid='ignore'):
        normalized = np.divide(cm, row_sums, where=row_sums != 0)
    normalized[np.isnan(normalized)] = 0.0
    return normalized

def save_confusion_matrix_images(cm, class_names, out_base, title_text):
    fig, ax = plt.subplots(figsize=(8, 6))
    im = ax.imshow(cm, interpolation='nearest', vmin=0.0, vmax=1.0, cmap=plt.cm.Blues)
    ax.set_title(title_text)
    plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    n = len(class_names)
    ax.set_xticks(np.arange(n))
    ax.set_yticks(np.arange(n))
    ax.set_xticklabels(class_names, rotation=45, ha="right")
    ax.set_yticklabels(class_names)
    thresh = 0.5
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(j, i, f"{cm[i, j]:.2f}",
                    ha="center", va="center",
                    color="white" if cm[i, j] > thresh else "black", fontsize=8)
    plt.tight_layout()
    paths = {}
    for ext in ("png", "svg", "pdf"):
        p = f"{out_base}.{ext}"
        fig.savefig(p, dpi=300, transparent=True)
        paths[ext] = p
    plt.close(fig)
    return paths

def add_slide_with_image(prs, image_path, slide_title):
    if Presentation is None:
        raise RuntimeError("python-pptx is required to create PPTX. Install with `pip install python-pptx`")
    layout_index = 6 if len(prs.slide_layouts) > 6 else 5
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(18)
    img_left = Inches(0.5)
    img_top = Inches(1.0)
    img_width = Inches(9.0)
    try:
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    except Exception as e:
        print(f"  Warning: could not insert image {image_path} into PPTX slide: {e}")

# ---------------- Main execution ----------------
def main():
    # quick sanity check
    if RESERVED_FINE_PCT + RF_VAL_PCT >= 1.0:
        raise ValueError("RESERVED_FINE_PCT + RF_VAL_PCT must be < 1.0 (not enough remaining data for RF training).")

    # Compute baseline embeddings if missing
    compute_missing_embeddings(device_list)

    # ---------------- Build single finetune set (first RESERVED_FINE_PCT per device) ----------------
    per_device_fine_idxs = {}
    all_fine_texts = []
    all_fine_labels = []
    for device in device_list:
        total = per_device_total_lines[device]
        fine_count = int(total * FINE_PCT)
        fine_count = max(1, fine_count)
        fine_idxs = list(range(0, fine_count))
        per_device_fine_idxs[device] = fine_idxs
        lines = per_device_lines[device]
        labels = [device_label_map[device]] * len(fine_idxs)
        per_device_texts = [lines[i] for i in fine_idxs]
        all_fine_texts.append(per_device_texts)
        all_fine_labels.append(labels)

    per_device_texts_map = {device_list[i]: all_fine_texts[i] for i in range(len(device_list))}
    per_device_labels_map = {device_list[i]: all_fine_labels[i] for i in range(len(device_list))}
    fine_texts, fine_labels = interleave_round_robin(per_device_texts_map, device_list, device_label_map)

    print(f"\nPrepared global fine-tune set from first {FINE_PCT*100:.1f}% per device: total fine examples = {len(fine_texts)}")

    # ---------------- Train RF teacher on baseline embeddings of the FIRST fine block ----------------
    print("\nTraining fixed RF teacher on baseline embeddings from the FIRST fine block (assumption)...")
    X_teacher = []
    y_teacher = []
    for device in device_list:
        emb_arr = load_embeddings_memmap_for_device(device, suffix="")
        if emb_arr is None:
            emb_arr = compute_and_cache_embeddings_stream(device, batch_size=EMBED_BATCH_SIZE)
        t_idxs = per_device_fine_idxs[device]
        if len(t_idxs) == 0:
            continue
        X_teacher.append(emb_arr[t_idxs, :])
        y_teacher.extend([device_label_map[device]] * len(t_idxs))
    if len(X_teacher) == 0:
        raise RuntimeError("No teacher training data found (unexpected).")
    X_teacher = np.vstack(X_teacher).astype(np.float32)
    y_teacher = np.array(y_teacher, dtype=int)
    print(f"  RF teacher training size: {X_teacher.shape[0]} examples")
    rf_teacher = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
    rf_teacher.fit(X_teacher, y_teacher)
    print("  RF teacher trained.")

    # ---------------- Decide whether to finetune or reuse existing finetuned embeddings ----------------
    skip_finetune = False

    if not FINETUNE_ENABLED:
        print("\nFINETUNE_ENABLED is False -> skipping finetuning and finetuned evaluation.")
        skip_finetune = True
    else:
        if REUSE_FINETUNE_EMBS_IF_PRESENT:
            ok = all_finetuned_embeddings_exist(device_list, suffix=".ft")
            if ok:
                print("\nAll finetuned embeddings (.ft.npy) already exist and REUSE_FINETUNE_EMBS_IF_PRESENT=True -> skipping finetuning.")
                skip_finetune = True
            else:
                print("\nFinetuned embeddings not present for all devices (or invalid); will perform finetuning and write .ft embeddings.")
        else:
            print("\nREUSE_FINETUNE_EMBS_IF_PRESENT=False -> (re)running finetuning regardless of existing .ft files.")

    # ---------------- STUDENT FINETUNE (single student trained once using fine_texts) ----------------
    student_model = None
    if not skip_finetune:
        print("\nStarting single student fine-tune using first-FINE_PCT interleaved examples (with RF feedback)...")
        # Prepare orig anchors and prototypes (orig embeddings over fine_texts) using cached baseline .npy where possible
        print("[2] Loading teacher (baseline) embeddings for the fine set from cache if present...")
        per_device_emb_slices = {}
        missing_cache = False

        # Try to read cached baseline .npy files and slice only the fine indices
        for device in device_list:
            emb_arr = load_embeddings_memmap_for_device(device, suffix="")
            if emb_arr is None:
                # missing or invalid cache -> fallback to computing embeddings for fine_texts below
                missing_cache = True
                break
            t_idxs = per_device_fine_idxs[device]
            if len(t_idxs) == 0:
                per_device_emb_slices[device] = []
            else:
                # slicing a memmap returns lightweight arrays per row; cast to float32 to be consistent
                per_device_emb_slices[device] = [np.array(emb_arr[i], dtype=np.float32) for i in t_idxs]

        if not missing_cache:
            # Interleave slices to match the fine_texts ordering used above
            orig_list, _ = interleave_round_robin(per_device_emb_slices, device_list, device_label_map)
            if len(orig_list) == 0:
                raise RuntimeError("No fine examples found when loading cached embeddings.")
            orig_cls_all = np.vstack(orig_list).astype(np.float32)
            print(f"[loaded] orig_cls_all shape: {orig_cls_all.shape}")
        else:
            # Fallback behavior: compute teacher embeddings only for the fine_texts (previous behavior),
            # rather than recomputing for everything.
            print("[fallback] Some baseline .npy files missing/invalid -> computing teacher embeddings for fine_texts now.")
            orig_cls_parts = []
            for i in range(0, len(fine_texts), FT_BATCH_SIZE):
                batch = fine_texts[i:i+FT_BATCH_SIZE]
                emb_t = model_batch_embeddings_torch(batch, tokenizer, base_model, max_length=BATCH_TOKEN_MAX_LEN)
                orig_cls_parts.append(emb_t.cpu().numpy())
            orig_cls_all = np.vstack(orig_cls_parts).astype(np.float32)
            print(f"[computed] orig_cls_all shape: {orig_cls_all.shape}")

        # Sanity check: ensure the number of rows matches the number of fine_texts
        if orig_cls_all.shape[0] != len(fine_texts):
            raise RuntimeError(f"orig_cls_all rows ({orig_cls_all.shape[0]}) != number of fine_texts ({len(fine_texts)}).")

        # Create labels array and centroids (unchanged)
        labels_arr = np.array(fine_labels, dtype=int)

        num_classes = len(device_list)
        centroids = np.zeros((num_classes, orig_cls_all.shape[1]), dtype=np.float32)
        for c in range(num_classes):
            idxs = np.where(labels_arr == c)[0]
            if len(idxs) > 0:
                centroids[c] = orig_cls_all[idxs].mean(axis=0)

        class FTData:
            def __init__(self, texts, labels, orig_embs, centroids):
                self.texts = texts
                self.labels = labels
                self.orig_embs = orig_embs.astype(np.float32)
                self.centroids = centroids
            def __len__(self):
                return len(self.texts)
            def __getitem__(self, idx):
                label = self.labels[idx]
                return self.texts[idx], int(label), self.orig_embs[idx], self.centroids[label]

        ft_dataset = FTData(fine_texts, fine_labels, orig_cls_all, centroids)
        from torch.utils.data import DataLoader
        ft_loader = DataLoader(ft_dataset, batch_size=FT_BATCH_SIZE, shuffle=False)

        # Student model + small classification head
        # For Mamba use AutoModelForCausalLM (same class as base_model); for encoder models use AutoModel
        if IS_MAMBA:
            student_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, torch_dtype=torch.float16).to(DEVICE)
        else:
            student_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg).to(DEVICE)

        hidden_size = student_model.config.hidden_size
        classifier_head = nn.Linear(hidden_size, num_classes).to(DEVICE)

        mse_loss = nn.MSELoss()
        ce_loss_per_sample = nn.CrossEntropyLoss(reduction='none')
        optimizer = torch.optim.AdamW(list(student_model.parameters()) + list(classifier_head.parameters()), lr=LR)

        print("  Fine-tune loop (anchor + proto + classification + RF-feedback)...")
        for epoch in range(EPOCHS):
            student_model.train()
            running_loss = 0.0
            steps = 0
            for batch in tqdm(ft_loader, desc=f" FT Epoch {epoch+1}/{EPOCHS}", leave=False):
                texts_b, labels_b, orig_emb_b, proto_emb_b = batch

                if isinstance(orig_emb_b, np.ndarray):
                    orig_for_rf = orig_emb_b
                else:
                    try:
                        orig_for_rf = orig_emb_b.cpu().numpy()
                    except Exception:
                        orig_for_rf = np.array(orig_emb_b)

                rf_preds = rf_teacher.predict(orig_for_rf)
                rf_preds = np.array(rf_preds, dtype=int)
                labels_np = np.array(labels_b, dtype=int)

                # Use unified helper to compute student embeddings (so Mamba/BERT consistent)
                enc_texts = list(texts_b)
                # compute student embeddings **with gradients enabled** so anchor/proto losses update the student
                cls_emb = model_batch_embeddings_torch(enc_texts, tokenizer, student_model, max_length=BATCH_TOKEN_MAX_LEN, require_grad=True)
                # ensure same dtype as orig_emb_t for stable MSE computation
                cls_emb = cls_emb.to(dtype=torch.float32, device=DEVICE)


                if isinstance(orig_emb_b, np.ndarray):
                    orig_emb_t = torch.from_numpy(orig_emb_b).to(dtype=torch.float32, device=DEVICE)
                else:
                    orig_emb_t = orig_emb_b.to(dtype=torch.float32, device=DEVICE)

                if isinstance(proto_emb_b, np.ndarray):
                    proto_emb_t = torch.from_numpy(proto_emb_b).to(dtype=torch.float32, device=DEVICE)
                else:
                    proto_emb_t = proto_emb_b.to(dtype=torch.float32, device=DEVICE)

                loss_anchor = mse_loss(cls_emb, orig_emb_t)
                loss_proto = mse_loss(cls_emb, proto_emb_t)

                logits = classifier_head(cls_emb)
                labels_t = torch.tensor(labels_np, dtype=torch.long, device=DEVICE)
                loss_cls_vec = ce_loss_per_sample(logits, labels_t)

                rf_pred_t = torch.tensor(rf_preds, dtype=torch.long, device=DEVICE)
                loss_rf_vec = ce_loss_per_sample(logits, rf_pred_t)
                signs = torch.where(rf_pred_t == labels_t, 1.0, -1.0).to(DEVICE)
                loss_rf = (loss_rf_vec * signs).mean()

                loss = (LAMBDA_ANCHOR * loss_anchor +
                        LAMBDA_PROTO * loss_proto +
                        LAMBDA_CLS * loss_cls_vec.mean() +
                        LAMBDA_RF  * loss_rf)
                optimizer.zero_grad()
                loss.backward()
                optimizer.step()

                running_loss += loss.item()
                steps += 1

            avg_loss = running_loss / max(1, steps)
            la = loss_anchor.item() if isinstance(loss_anchor, torch.Tensor) else float(loss_anchor)
            lp = loss_proto.item() if isinstance(loss_proto, torch.Tensor) else float(loss_proto)
            lc = loss_cls_vec.mean().item() if isinstance(loss_cls_vec, torch.Tensor) else float(loss_cls_vec.mean())
            lrf = loss_rf.item() if isinstance(loss_rf, torch.Tensor) else float(loss_rf)
            print(f"   Epoch {epoch+1}/{EPOCHS} - avg loss: {avg_loss:.6f} (anchor {la:.6f}, proto {lp:.6f}, cls {lc:.6f}, rf {lrf:.6f})")

        print("  Finished finetuning student. Saving finetuned embeddings (.ft) per device...")

        for device in device_list:
            compute_and_cache_embeddings_stream_for_model(device, student_model, suffix=".ft", batch_size=EMBED_BATCH_SIZE)

        print("  Finetuned embeddings saved.")
    else:
        print("\nSkipping finetuning - will use existing .ft embeddings (if present) or skip finetuned evaluation entirely.")

    # ---------------- RF sweep training/eval using baseline and finetuned embeddings ----------------
    results_baseline = []
    results_finetuned = []
    cm_baseline_records = []
    cm_finetuned_records = []
    classes = [d.replace('.json', '') for d in device_list]

    rf_train_pcts = np.arange(START_PCT, END_PCT, STEP_PCT)

    print(f"Reserved fine block: first {RESERVED_FINE_PCT*100:.1f}% per device (always excluded from RF training).")
    if START_PCT <= RESERVED_FINE_PCT:
        print(f"Note: START_PCT = {START_PCT*100:.1f}% <= RESERVED_FINE_PCT = {RESERVED_FINE_PCT*100:.1f}%. The initial sweep step(s) may have empty RF training sets until RF_TRAIN_PCT > {RESERVED_FINE_PCT*100:.1f}%.")

    for pct in rf_train_pcts:
        RF_TRAIN_PCT = float(pct)
        print("------------------------------------------------------------")
        print(f"Running sweep step: RF_TRAIN_PCT = {RF_TRAIN_PCT*100:.1f}% (RF-val fixed at last {RF_VAL_PCT*100:.1f}%)")
        per_device_train_emb_idxs = {}
        per_device_val_emb_idxs = {}
        total_train = 0
        total_val = 0
        for device in device_list:
            total_lines = per_device_total_lines[device]
            rf_val_count = int(total_lines * RF_VAL_PCT)
            rf_val_count = max(1, rf_val_count)

            # use RESERVED_FINE_PCT so reserved/fine block is always the fixed first 20%
            fine_end_idx = int(total_lines * RESERVED_FINE_PCT)
            fine_end_idx = min(fine_end_idx, total_lines - rf_val_count)
            if fine_end_idx < 0:
                fine_end_idx = 0

            rf_train_start_idx = fine_end_idx
            rf_train_end_idx = int(total_lines * RF_TRAIN_PCT)
            rf_train_end_idx = min(rf_train_end_idx, total_lines - rf_val_count)
            if rf_train_end_idx <= rf_train_start_idx:
                train_idxs = []
            else:
                train_idxs = list(range(rf_train_start_idx, rf_train_end_idx))

            val_idxs = list(range(total_lines - rf_val_count, total_lines))

            per_device_train_emb_idxs[device] = train_idxs
            per_device_val_emb_idxs[device] = val_idxs
            total_train += len(train_idxs)
            total_val += len(val_idxs)

            if PRINT_SAMPLE_LINES:
                lines = per_device_lines[device]
                t_first = short(lines[train_idxs[0]]) if train_idxs else "<none>"
                t_last = short(lines[train_idxs[-1]]) if train_idxs else "<none>"
                v_first = short(lines[val_idxs[0]]) if val_idxs else "<none>"
                v_last = short(lines[val_idxs[-1]]) if val_idxs else "<none>"
                fine_first = short(lines[0]) if fine_end_idx > 0 else "<none>"
                fine_last = short(lines[fine_end_idx-1]) if fine_end_idx > 0 else "<none>"
                print(f"  {device}: total={total_lines}, fine 1-based=(1,{fine_end_idx}) count={fine_end_idx}, train 1-based=({rf_train_start_idx+1},{rf_train_end_idx}) count={len(train_idxs)}, val 1-based=({total_lines-len(val_idxs)+1},{total_lines}) count={len(val_idxs)}")
                if fine_end_idx > 0:
                    print(f"    fine sample first/last: {fine_first} / {fine_last}")
                if train_idxs:
                    print(f"    train sample first/last: {t_first} / {t_last}")
                if val_idxs:
                    print(f"    val   sample first/last: {v_first} / {v_last}")

        print(f"  Combined totals before interleave: n_train={total_train}, n_val={total_val}")

        # Baseline
        per_device_train_lists = {}
        per_device_val_lists = {}
        for device in device_list:
            emb_arr = load_embeddings_memmap_for_device(device, suffix="")
            if emb_arr is None:
                emb_arr = compute_and_cache_embeddings_stream(device, batch_size=EMBED_BATCH_SIZE)
            t_idxs = per_device_train_emb_idxs[device]
            v_idxs = per_device_val_emb_idxs[device]
            per_device_train_lists[device] = [np.array(emb_arr[i], dtype=np.float32) for i in t_idxs] if len(t_idxs) > 0 else []
            per_device_val_lists[device]   = [np.array(emb_arr[i], dtype=np.float32) for i in v_idxs] if len(v_idxs) > 0 else []

        train_emb_list, train_labels = interleave_round_robin(per_device_train_lists, device_list, device_label_map)
        val_emb_list, val_labels     = interleave_round_robin(per_device_val_lists, device_list, device_label_map)

        if len(train_emb_list) == 0:
            print("  No baseline training samples available for this pct; skipping baseline.")
        else:
            X_train = np.vstack(train_emb_list).astype(np.float32)
            y_train = np.array(train_labels, dtype=int)
            X_val = np.vstack(val_emb_list).astype(np.float32) if len(val_emb_list) > 0 else np.zeros((0, X_train.shape[1]), dtype=np.float32)
            y_val = np.array(val_labels, dtype=int) if len(val_emb_list) > 0 else np.array([], dtype=int)

            print(f"  Baseline combined totals after interleave: n_train={X_train.shape[0]}, n_val={X_val.shape[0]}")
            rf = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
            rf.fit(X_train, y_train)
            if X_val.shape[0] == 0:
                print("  No baseline validation samples; skipping baseline eval.")
            else:
                y_pred = rf.predict(X_val)
                f1m = f1_score(y_val, y_pred, average='macro')
                print(f"  -> Baseline Macro F1 at RF_TRAIN_PCT={RF_TRAIN_PCT*100:.1f}%: {f1m:.4f}")
                cm = confusion_matrix(y_val, y_pred, labels=list(range(len(device_list))))
                cm_norm = normalize_cm_rows(cm)
                out_base = os.path.join(CM_BASELINE_DIR, f"confusion_baseline_RFtrain_{RF_TRAIN_PCT*100:.1f}pct")
                title_text = f"Baseline RF Train {RF_TRAIN_PCT*100:.1f}% — Macro F1: {f1m:.4f}"
                saved = save_confusion_matrix_images(cm_norm, classes, out_base, title_text)
                cm_baseline_records.append({"pct": RF_TRAIN_PCT, "macro_f1": float(f1m), "n_train": int(X_train.shape[0]), "n_val": int(X_val.shape[0]), "img_png": saved.get("png"), "cm": cm_norm})
                results_baseline.append({"rf_train_pct": RF_TRAIN_PCT, "macro_f1": float(f1m), "n_train": int(X_train.shape[0]), "n_val": int(X_val.shape[0])})

        # Finetuned (only if FINETUNE_ENABLED)
        if FINETUNE_ENABLED:
            per_device_train_lists_ft = {}
            per_device_val_lists_ft = {}
            for device in device_list:
                emb_arr_ft = load_embeddings_memmap_for_device(device, suffix=".ft")
                if emb_arr_ft is None:
                    if student_model is None:
                        # fallback: create an unfinetuned student model to compute embeddings
                        if IS_MAMBA:
                            student_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, torch_dtype=torch.float16).to(DEVICE)
                        else:
                            student_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg).to(DEVICE)
                    emb_arr_ft = compute_and_cache_embeddings_stream_for_model(device, student_model, suffix=".ft", batch_size=EMBED_BATCH_SIZE)
                t_idxs = per_device_train_emb_idxs[device]
                v_idxs = per_device_val_emb_idxs[device]
                per_device_train_lists_ft[device] = [np.array(emb_arr_ft[i], dtype=np.float32) for i in t_idxs] if len(t_idxs) > 0 else []
                per_device_val_lists_ft[device]   = [np.array(emb_arr_ft[i], dtype=np.float32) for i in v_idxs] if len(v_idxs) > 0 else []

            train_emb_list_ft, train_labels_ft = interleave_round_robin(per_device_train_lists_ft, device_list, device_label_map)
            val_emb_list_ft, val_labels_ft     = interleave_round_robin(per_device_val_lists_ft, device_list, device_label_map)

            if len(train_emb_list_ft) == 0:
                print("  No finetuned training samples available for this pct; skipping finetuned RF.")
            else:
                X_train_ft = np.vstack(train_emb_list_ft).astype(np.float32)
                y_train_ft = np.array(train_labels_ft, dtype=int)
                X_val_ft = np.vstack(val_emb_list_ft).astype(np.float32) if len(val_emb_list_ft) > 0 else np.zeros((0, X_train_ft.shape[1]), dtype=np.float32)
                y_val_ft = np.array(val_labels_ft, dtype=int) if len(val_emb_list_ft) > 0 else np.array([], dtype=int)

                print(f"  Finetuned combined totals after interleave: n_train={X_train_ft.shape[0]}, n_val={X_val_ft.shape[0]}")
                rf_ft = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
                rf_ft.fit(X_train_ft, y_train_ft)
                if X_val_ft.shape[0] == 0:
                    print("  No finetuned validation samples; skipping finetuned eval.")
                else:
                    y_pred_ft = rf_ft.predict(X_val_ft)
                    f1m_ft = f1_score(y_val_ft, y_pred_ft, average='macro')
                    print(f"  -> Finetuned Macro F1 at RF_TRAIN_PCT={RF_TRAIN_PCT*100:.1f}%: {f1m_ft:.4f}")
                    cm_ft = confusion_matrix(y_val_ft, y_pred_ft, labels=list(range(len(device_list))))
                    cm_ft_norm = normalize_cm_rows(cm_ft)
                    out_base_ft = os.path.join(CM_FINETUNED_DIR, f"confusion_finetuned_RFtrain_{int(RF_TRAIN_PCT*100):02d}pct")
                    title_text_ft = f"Finetuned RF Train {RF_TRAIN_PCT*100:.1f}% — Macro F1: {f1m_ft:.4f}"
                    saved_ft = save_confusion_matrix_images(cm_ft_norm, classes, out_base_ft, title_text_ft)
                    cm_finetuned_records.append({"pct": RF_TRAIN_PCT, "macro_f1": float(f1m_ft), "n_train": int(X_train_ft.shape[0]), "n_val": int(X_val_ft.shape[0]), "img_png": saved_ft.get("png"), "cm": cm_ft_norm})
                    results_finetuned.append({"rf_train_pct": RF_TRAIN_PCT, "macro_f1": float(f1m_ft), "n_train": int(X_train_ft.shape[0]), "n_val": int(X_val_ft.shape[0])})
        else:
            print("  Finetuned evaluation skipped (FINETUNE_ENABLED=False).")

    # ---------------- Save numeric results and plots ----------------
    def save_results_csv_and_plot(results_list, csv_name, plot_name):
        if len(results_list) == 0:
            return None
        csv_path = csv_name
        with open(csv_path, "w", newline='', encoding='utf-8') as csvf:
            writer = csv.DictWriter(csvf, fieldnames=["rf_train_pct", "n_train", "n_val", "macro_f1"])
            writer.writeheader()
            for r in results_list:
                writer.writerow({"rf_train_pct": r["rf_train_pct"], "n_train": r["n_train"], "n_val": r["n_val"], "macro_f1": r["macro_f1"]})
        pcts = [r["rf_train_pct"] * 100.0 for r in results_list]
        f1s = [r["macro_f1"] for r in results_list]
        plt.figure(figsize=(8, 5))
        plt.plot(pcts, f1s, marker='o')
        plt.xlabel("RF_TRAIN_PCT (%)")
        plt.ylabel("Macro F1 (validation)")
        plt.title(plot_name)
        plt.grid(True)
        plt.tight_layout()
        plot_path = csv_name.replace(".csv", ".png")
        plt.savefig(plot_path, dpi=300, transparent=True)
        plt.close()
        return csv_path, plot_path

    base_csv, base_plot = None, None
    if results_baseline:
        base_csv, base_plot = save_results_csv_and_plot(results_baseline, "rf_train_pct_results_baseline.csv",
                                                        "Baseline: RF training size vs Macro F1")
        print(f"\nSaved baseline numeric results to: {base_csv}")

    ft_csv, ft_plot = None, None
    if FINETUNE_ENABLED and results_finetuned:
        ft_csv, ft_plot = save_results_csv_and_plot(results_finetuned, "rf_train_pct_results_finetuned.csv",
                                                    "Finetuned: RF training size vs Macro F1")
        print(f"\nSaved finetuned numeric results to: {ft_csv}")
    elif not FINETUNE_ENABLED:
        print("\nFinetuning disabled -> no finetuned results saved.")

    # ---------------- Create PPTX containing confusion matrices ----------------
    def create_pptx_from_records(records, out_pptx):
        if Presentation is None:
            print(f"Skipping PPTX creation ({out_pptx}): python-pptx not installed.")
            return
        prs = Presentation()
        try:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = "Confusion Matrices"
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated: {datetime.now(ZoneInfo('Australia/Sydney')).isoformat()} Sydney time"

        except Exception:
            pass
        for rec in records:
            title = f"RF Train {rec['pct']*100:.1f}% — Macro F1: {rec['macro_f1']:.4f} — n_train={rec['n_train']} n_val={rec['n_val']}"
            img = rec.get("img_png")
            if img is None:
                layout_index = 1 if len(prs.slide_layouts) > 1 else 0
                s = prs.slides.add_slide(prs.slide_layouts[layout_index])
                tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
                tb.text_frame.text = title
                continue
            add_slide_with_image(prs, img, title)
        prs.save(out_pptx)
        print(f"Saved PPTX: {out_pptx}")

    pptx_base_path = None
    pptx_ft_path = None
    if cm_baseline_records:
        pptx_base_path = os.path.join(CM_DIR, f"confusion_baseline_{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}.pptx")
        create_pptx_from_records(cm_baseline_records, pptx_base_path)

    if FINETUNE_ENABLED and cm_finetuned_records:
        pptx_ft_path = os.path.join(CM_DIR, f"confusion_finetuned_{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}.pptx")
        create_pptx_from_records(cm_finetuned_records, pptx_ft_path)

    print("\nDone. Files generated:")
    if base_csv:
        print(f"  - baseline CSV: {base_csv}")
    if base_plot:
        print(f"  - baseline plot: {base_plot}")
    if FINETUNE_ENABLED and ft_csv:
        print(f"  - finetuned CSV: {ft_csv}")
    if FINETUNE_ENABLED and ft_plot:
        print(f"  - finetuned plot: {ft_plot}")
    print(f"  - confusion matrices (baseline) in: {CM_BASELINE_DIR}")
    if FINETUNE_ENABLED:
        print(f"  - confusion matrices (finetuned) in: {CM_FINETUNED_DIR}")
    else:
        print("  - finetuned outputs skipped (FINETUNE_ENABLED=False)")
    print(f"  - embedding files (per device) in folder: {EMB_DIR}")

    # Return a summary for the success email
    summary = {
        "baseline_csv": base_csv,
        "baseline_plot": base_plot,
        "finetuned_csv": ft_csv if FINETUNE_ENABLED else None,
        "finetuned_plot": ft_plot if FINETUNE_ENABLED else None,
        "pptx_baseline": pptx_base_path,
        "pptx_finetuned": pptx_ft_path if FINETUNE_ENABLED else None,
        "n_baseline_results": len(results_baseline),
        "n_finetuned_results": len(results_finetuned) if FINETUNE_ENABLED else 0,
        "finetune_enabled": FINETUNE_ENABLED,
    }
    return summary

# ---------------- Run main inside wrapper and send emails on success/failure ----------------
if __name__ == "__main__":
    try:
        summary = main()
        # success email
        subject = "Code Completed Successfully"
        body_lines = [
            "Your RF sweep + finetune script completed successfully.",
            "",
            f"Model: {MODEL_NAME}",
            f"Finetune enabled: {summary.get('finetune_enabled', False)}",
            f"Baseline results points: {summary.get('n_baseline_results', 0)}",
            f"Finetuned results points: {summary.get('n_finetuned_results', 0)}",
            ""
        ]
        if summary.get("baseline_csv"):
            body_lines.append(f"Baseline CSV: {summary['baseline_csv']}")
        if summary.get("baseline_plot"):
            body_lines.append(f"Baseline plot: {summary['baseline_plot']}")
        if summary.get("finetuned_csv"):
            body_lines.append(f"Finetuned CSV: {summary['finetuned_csv']}")
        if summary.get("finetuned_plot"):
            body_lines.append(f"Finetuned plot: {summary['finetuned_plot']}")
        if summary.get("pptx_baseline"):
            body_lines.append(f"Baseline PPTX: {summary['pptx_baseline']}")
        if summary.get("pptx_finetuned"):
            body_lines.append(f"Finetuned PPTX: {summary['pptx_finetuned']}")
        body = "\n".join(body_lines)
        try:
            special.send_test_email(subject, body)
        except Exception as send_e:
            print(f"Warning: sending success email failed: {send_e}")
            print("Success email body would have been:\n", body)
    except Exception as e:
        error_subject = "Code Failed!"
        tb_str = traceback.format_exc()  # Get the full traceback as a string
        error_body = f"The code encountered an error:\n{str(e)}\n\nFull traceback:\n{tb_str}"
        try:
            special.send_test_email(error_subject, error_body)
        except Exception as send_e:
            print(f"Warning: sending error email failed: {send_e}")
            print("Error body would have been:\n", error_body)
        # re-raise so the error is still visible on console/logs
        raise
