#!/usr/bin/env python3
# new_pipeline_gpu_fp32_parity.py
"""
GPU script tuned to PRODUCE THE SAME EMBEDDINGS as CPU:
 - Forces fp32 arithmetic everywhere (no torch_dtype=float16)
 - Removes autocast and GradScaler for embedding forward passes and training loops
 - Uses CUDA device if available (for speed) but computations remain float32 to match CPU
 - Keeps your OOM fallback, caching, and RF sweep logic intact

Minimal edits:
 - use torch.inference_mode() for inference helper
 - vectorized memmap writes (bulk slice assignment)
 - temporarily disable gradient checkpointing only during .ft streaming inference (re-enabled afterwards)
 - free memory and delete temporaries inside streaming loops

This version adds a small diagnostic block (after saving a finetuned model)
that compares embeddings from the in-memory student model vs a reloaded
saved copy under a few settings to help debug discrepancies.
"""
import os
import csv
import textwrap
import numpy as np
from tqdm import tqdm
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import confusion_matrix, f1_score
from transformers import AutoTokenizer, AutoModel, AutoConfig, AutoModelForCausalLM
import torch
import torch.nn as nn
import sys
import matplotlib.pyplot as plt
from datetime import datetime
from torch.utils.data import DataLoader
import time
import matplotlib.pyplot as plt

# email & traceback imports for wrapper
import traceback
import special

# try to import python-pptx; if missing we'll raise a clear error when needed
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None




# ---------------- Determinism / reproducibility (best-effort) ----------------
# Set seeds and deterministic flags to reduce non-determinism
# os.environ["CUDA_LAUNCH_BLOCKING"] = "1"
os.environ["PYTHONHASHSEED"] = "0"
torch.manual_seed(42)
if torch.cuda.is_available():
    torch.cuda.manual_seed_all(42)
torch.backends.cudnn.deterministic = True
torch.backends.cudnn.benchmark = False

# ---------------- Config ----------------
device_list = [
    "irobot_roomba.json",
    "line_clova_wave.json",
    "nature_remo.json",
    "qrio_hub.json",
    "xiaomi_mijia_led.json",
    "powerelectric_wi-fi_plug.json",
    "planex_smacam_outdoor.json"
]
# adjust paths as needed
INPUT_FOLDER = "preprocessed_data_group_merged/ungrouped"

# Model selection
MODEL_NAME = "state-spaces/mamba-130m-hf"

# Detect Mamba-like model (minimal rule)
IS_MAMBA = ("mamba" in MODEL_NAME.lower()) or MODEL_NAME.startswith("state-spaces/mamba")

# Fraction of each device file to use (0.0 < DATA_USAGE_PCT <= 1.0)
# Set to 0.5 to use the first 50% of each device's data.
DATA_USAGE_PCT = 0.1  # default 1.0 (use everything); change to 0.5, 0.25, etc.

# Finetune flag
FINETUNE_ENABLED = False
DEBUG_FT_VERBOSE = True
student_model = None

# Sweep settings
START_PCT = 0.20
END_PCT = 0.70
STEP_PCT = 0.02
RF_VAL_PCT = 0.30
RESERVED_FINE_PCT = 0.20

# Fine-tune specification
# FINE_PCT = 0.01
FINE_PCT = 0.0001

REUSE_FINETUNE_EMBS_IF_PRESENT = True

# DEVICE (cuda if available)
DEVICE = torch.device("cuda" if torch.cuda.is_available() else "cpu")
IS_CUDA = DEVICE.type == "cuda"
if IS_CUDA:
    try:
        print(f"Using CUDA: {torch.cuda.get_device_name(0)}")
    except Exception:
        print("Using CUDA")
else:
    print("CUDA not available; running on CPU")

# ---------------- Auto-tune batch sizes (same logic as before) ----------------
def bytes_to_gb(x): return float(x) / (1024**3)

BATCH_TOKEN_MAX_LEN = 512
EMBED_BATCH_SIZE = 12
FT_BATCH_SIZE = 1

# if IS_CUDA:
#     try:
#         props = torch.cuda.get_device_properties(0)
#         total_gb = bytes_to_gb(props.total_memory)
#         if total_gb < 4.5:
#             EMBED_BATCH_SIZE = 8
#             # EMBED_BATCH_SIZE = 2
#             EMBED_BATCH_SIZE = 4
#         elif total_gb < 8:
#             EMBED_BATCH_SIZE = 16
#         elif total_gb < 16:
#             EMBED_BATCH_SIZE = 32
#         else:
#             EMBED_BATCH_SIZE = 64
#     except Exception:
#         pass
# else:
#     EMBED_BATCH_SIZE = 32

# OOM-retry wrapper (same as your GPU script)
def run_with_oom_fallback(func, *, batch_arg_name="batch_size", start_batch=None, min_batch=1, **kwargs):
    current = start_batch if start_batch is not None else (EMBED_BATCH_SIZE if batch_arg_name == "batch_size" else FT_BATCH_SIZE)
    while True:
        try:
            kwargs[batch_arg_name] = int(max(min_batch, current))
            if IS_CUDA:
                torch.cuda.empty_cache()
            return func(**kwargs)
        except RuntimeError as e:
            m = str(e).lower()
            if "out of memory" in m or "cuda out of memory" in m:
                new = max(min_batch, current // 2)
                if new >= current:
                    raise
                print(f"CUDA OOM with batch {current} -> retrying with batch {new}")
                current = new
                if IS_CUDA:
                    torch.cuda.empty_cache()
            else:
                raise

RF_N_ESTIMATORS = 500
RF_N_JOBS = -1

# Finetuning hyperparams
LR = 2e-5
EPOCHS = 3
LAMBDA_ANCHOR = 0.5
LAMBDA_PROTO = 0.5
LAMBDA_CLS = 0.5
LAMBDA_RF  = 0.5

# Embedding cache folder
EMB_DIR = "embeddings"
FINETUNED_MODEL_DIR = os.path.join(EMB_DIR, "finetuned_model")
CM_DIR = os.path.join(EMB_DIR, "confusion_matrices")
CM_BASELINE_DIR = os.path.join(CM_DIR, "baseline")
CM_FINETUNED_DIR = os.path.join(CM_DIR, "finetuned")
os.makedirs(EMB_DIR, exist_ok=True)
os.makedirs(CM_DIR, exist_ok=True)
os.makedirs(CM_BASELINE_DIR, exist_ok=True)
os.makedirs(CM_FINETUNED_DIR, exist_ok=True)

PRINT_SAMPLE_LINES = True

# ---------------- Helpers ----------------
def short(s, width=120):
    return textwrap.shorten(s.replace("\n", "\\n"), width=width, placeholder="…")

def count_nonempty_lines(path):
    cnt = 0
    with open(path, 'r', encoding='utf-8') as f:
        for line in f:
            if line.strip():
                cnt += 1
    return cnt

def load_nonempty_lines(file_path):
    lines = []
    with open(file_path, 'r', encoding='utf-8') as f:
        for line in f:
            s = line.rstrip("\n")
            if s.strip():
                lines.append(s)
    return lines

# ---------------- Prepare environment (model/tokenizer) ----------------
print("Loading tokenizer and model (this may take a moment)...")
tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME, use_fast=True)
cfg = AutoConfig.from_pretrained(MODEL_NAME)
cfg.output_hidden_states = True

# IMPORTANT: no torch_dtype=float16 — load in fp32 for parity with CPU
dtype_kwargs = {}

# Minimal branching for Mamba vs encoder models:
if IS_MAMBA:
    base_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)
else:
    base_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)

# Move model to device (parameters remain fp32)
base_model.to(DEVICE)
base_model.eval()

try:
    base_model.gradient_checkpointing_enable()
    print("[info] Enabled gradient checkpointing on base_model (saves activations).", flush=True)
except Exception:
    pass

device_label_map = {dev: idx for idx, dev in enumerate(device_list)}
device_short_names = [d.replace('.json', '') for d in device_list]

# # ---------------- Precompute / load lines and set up caches ----------------
# per_device_lines = {}
# per_device_total_lines = {}

# for device in device_list:
#     file_name = device.replace('.json', '.json.txt')
#     file_path = os.path.join(INPUT_FOLDER, file_name)
#     if not os.path.exists(file_path):
#         raise FileNotFoundError(f"Missing file: {file_path}")

#     total_lines = count_nonempty_lines(file_path)
#     if total_lines == 0:
#         raise ValueError(f"No non-empty lines in {file_path}")
#     per_device_total_lines[device] = total_lines

#     safe_name = device.replace('.json', '')
#     lines_cache_path = os.path.join(EMB_DIR, f"{safe_name}.lines.txt")

#     if os.path.exists(lines_cache_path):
#         lines = load_nonempty_lines(lines_cache_path)
#         if len(lines) != total_lines:
#             lines = load_nonempty_lines(file_path)
#             with open(lines_cache_path, 'w', encoding='utf-8') as wf:
#                 for L in lines:
#                     wf.write(L + "\n")
#     else:
#         lines = load_nonempty_lines(file_path)
#         with open(lines_cache_path, 'w', encoding='utf-8') as wf:
#             for L in lines:
#                 wf.write(L + "\n")

#     per_device_lines[device] = lines

# ---------------- Precompute / load lines and set up caches (with optional truncation) ----------------
per_device_lines = {}
per_device_total_lines = {}

for device in device_list:
    file_name = device.replace('.json', '.json.txt')
    file_path = os.path.join(INPUT_FOLDER, file_name)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Missing file: {file_path}")

    total_lines = count_nonempty_lines(file_path)
    if total_lines == 0:
        raise ValueError(f"No non-empty lines in {file_path}")

    # Compute how many lines we will actually use (truncate to the first fraction)
    use_count = int(total_lines * DATA_USAGE_PCT)
    use_count = max(1, min(total_lines, use_count))  # at least 1, at most total_lines
    per_device_total_lines[device] = use_count

    safe_name = device.replace('.json', '')
    lines_cache_path = os.path.join(EMB_DIR, f"{safe_name}.lines.txt")

    # If cache exists and matches requested truncated length, reuse it.
    # Otherwise, read original file and write truncated cache.
    if os.path.exists(lines_cache_path):
        lines = load_nonempty_lines(lines_cache_path)
        # If cache length != desired use_count, regenerate truncated cache from original file
        if len(lines) != use_count:
            full_lines = load_nonempty_lines(file_path)
            lines = full_lines[:use_count]
            with open(lines_cache_path, 'w', encoding='utf-8') as wf:
                for L in lines:
                    wf.write(L + "\n")
    else:
        full_lines = load_nonempty_lines(file_path)
        lines = full_lines[:use_count]
        with open(lines_cache_path, 'w', encoding='utf-8') as wf:
            for L in lines:
                wf.write(L + "\n")

    per_device_lines[device] = lines


print("\nAll device line counts:")
sum_device_count = 0
for d in device_list:
    print(f"  {d}: total non-empty lines = {per_device_total_lines[d]}")
    sum_device_count += per_device_total_lines[d]
print(f"Total lines across all devices: {sum_device_count}")
print("")

# ---------------- Unified embedding helper (fp32, no autocast) ----------------
def model_batch_embeddings_torch(batch_texts, tokenizer, model, max_length=BATCH_TOKEN_MAX_LEN):
    """
    Returns torch.Tensor shape (B, H) on DEVICE in float32.
    Uses torch.inference_mode() for fastest inference where appropriate.
    """
    enc = tokenizer(batch_texts, return_tensors='pt', truncation=True, max_length=max_length, padding=True)
    enc = {k: v.to(DEVICE) for k, v in enc.items()}

    # prefer inference_mode when available (faster than no_grad)
    inference_ctx = torch.inference_mode if hasattr(torch, "inference_mode") else torch.no_grad
    with inference_ctx():
        out = model(**enc)

    if IS_MAMBA:
        last_hidden = out.hidden_states[-1]  # (B, seq_len, hidden)
        mask = enc["attention_mask"].unsqueeze(-1).to(last_hidden.dtype)  # (B, seq_len, 1)
        summed = (last_hidden * mask).sum(dim=1)
        counts = mask.sum(dim=1).clamp(min=1)
        emb = summed / counts
    else:
        emb = out.last_hidden_state[:, 0, :]

    # ensure float32 dtype and on DEVICE
    return emb.to(dtype=torch.float32, device=DEVICE)

def get_embedding_batch(sentences, tokenizer, model, max_length=BATCH_TOKEN_MAX_LEN):
    emb_t = model_batch_embeddings_torch(sentences, tokenizer, model, max_length=max_length)
    return emb_t.cpu().numpy()

def get_embedding_batch_on_device(sentences, tokenizer, model, device, max_length=BATCH_TOKEN_MAX_LEN):
    """
    Like model_batch_embeddings_torch but uses a specific device (useful for CPU comparisons).
    Returns numpy (B, H) on CPU.
    """
    enc = tokenizer(sentences, return_tensors='pt', truncation=True, max_length=max_length, padding=True)
    enc = {k: v.to(device) for k, v in enc.items()}
    inference_ctx = torch.inference_mode if hasattr(torch, "inference_mode") else torch.no_grad
    with inference_ctx():
        out = model(**enc)
    if IS_MAMBA:
        last_hidden = out.hidden_states[-1]
        mask = enc["attention_mask"].unsqueeze(-1).to(last_hidden.dtype)
        summed = (last_hidden * mask).sum(dim=1)
        counts = mask.sum(dim=1).clamp(min=1)
        emb = summed / counts
    else:
        emb = out.last_hidden_state[:, 0, :]
    return emb.to(dtype=torch.float32, device=device).cpu().numpy()

def interleave_round_robin(per_device_lists, devices_order, device_label_map):
    pointers = {d: 0 for d in devices_order}
    total_remaining = sum(len(per_device_lists[d]) for d in devices_order)
    combined = []
    labels = []
    while total_remaining > 0:
        for d in devices_order:
            idx = pointers[d]
            lst = per_device_lists[d]
            if idx < len(lst):
                combined.append(lst[idx])
                labels.append(device_label_map[d])
                pointers[d] += 1
                total_remaining -= 1
    return combined, labels

# ---------------- Streaming embedding computation & caching ----------------
from numpy.lib.format import open_memmap

def compute_and_cache_embeddings_stream(device, batch_size=EMBED_BATCH_SIZE):
    safe_name = device.replace('.json', '')
    lines = per_device_lines[device]
    n = len(lines)
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb.npy")
    emb_txt_path = os.path.join(EMB_DIR, f"{safe_name}.emb.txt")

    if os.path.exists(emb_npy_path):
        try:
            candidate = np.load(emb_npy_path, mmap_mode='r')
            if candidate.ndim == 2 and candidate.shape[0] == n and candidate.shape[1] == base_model.config.hidden_size:
                print(f"  Existing .npy for {device} looks valid; reusing: {emb_npy_path}")
                return candidate
            else:
                try:
                    os.remove(emb_npy_path)
                except OSError:
                    pass
                if os.path.exists(emb_txt_path):
                    try:
                        os.remove(emb_txt_path)
                    except OSError:
                        pass
        except Exception:
            try:
                os.remove(emb_npy_path)
            except OSError:
                pass
            if os.path.exists(emb_txt_path):
                try:
                    os.remove(emb_txt_path)
                except OSError:
                    pass

    hidden_size = base_model.config.hidden_size
    print(f"Streaming-computing embeddings for {device} ({n} lines). Writing .npy -> {emb_npy_path}")
    mm = open_memmap(emb_npy_path, dtype='float32', mode='w+', shape=(n, hidden_size))
    txt_f = open(emb_txt_path, 'w', encoding='utf-8')

    base_model.eval()
    # use inference_mode for faster inference
    inference_ctx = torch.inference_mode if hasattr(torch, "inference_mode") else torch.no_grad
    with inference_ctx():
        idx = 0
        for start in tqdm(range(0, n, batch_size), desc=f"Embedding (stream) {safe_name}"):
            batch_texts = lines[start:start+batch_size]
            batch_embs = get_embedding_batch(batch_texts, tokenizer, base_model, max_length=BATCH_TOKEN_MAX_LEN)
            B = batch_embs.shape[0]

            # bulk write to memmap slice
            mm[idx:idx+B, :] = batch_embs.astype(np.float32)

            # write text lines for this batch
            for b in range(B):
                row = batch_embs[b]
                txt_f.write(" ".join(f"{float(v):.6e}" for v in row) + "\n")

            idx += B
            # free intermediates and reduce fragmentation
            del batch_embs
            if IS_CUDA:
                torch.cuda.empty_cache()

    txt_f.flush()
    txt_f.close()
    mm.flush()
    arr = np.load(emb_npy_path, mmap_mode='r')
    print(f"  Finished {device}. .npy saved at {emb_npy_path} shape={arr.shape}")
    return arr

def compute_and_cache_embeddings_stream_for_model(device, model, suffix=".ft", batch_size=EMBED_BATCH_SIZE):
    """
    Compute embeddings for a model instance (e.g. the finetuned student model).
    To reduce OOM and speed: temporarily disable checkpointing only during this inference pass (then restore).
    Also uses vectorized memmap writes and inference_mode.
    """
    safe_name = device.replace('.json', '')
    lines = per_device_lines[device]
    n = len(lines)
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.npy")
    emb_txt_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.txt")

    if os.path.exists(emb_npy_path):
        try:
            candidate = np.load(emb_npy_path, mmap_mode='r')
            if candidate.ndim == 2 and candidate.shape[0] == n and candidate.shape[1] == model.config.hidden_size:
                print(f"  Existing .npy (suffix={suffix}) for {device} looks valid; reusing: {emb_npy_path}")
                return candidate
            else:
                try:
                    os.remove(emb_npy_path)
                except OSError:
                    pass
                if os.path.exists(emb_txt_path):
                    try:
                        os.remove(emb_txt_path)
                    except OSError:
                        pass
        except Exception:
            try:
                os.remove(emb_npy_path)
            except OSError:
                pass
            if os.path.exists(emb_txt_path):
                try:
                    os.remove(emb_txt_path)
                except OSError:
                    pass

    hidden_size = model.config.hidden_size
    print(f"Streaming-computing embeddings (model suffix={suffix}) for {device} ({n} lines). Writing .npy -> {emb_npy_path}")
    mm = open_memmap(emb_npy_path, dtype='float32', mode='w+', shape=(n, hidden_size))
    txt_f = open(emb_txt_path, 'w', encoding='utf-8')

    # We'll perform inference in inference_mode and temporarily disable checkpointing
    # Note: we DO NOT touch checkpointing during finetuning; this only runs AFTER training
    checkpoint_was_enabled = False
    try:
        # detect and disable if available (safe to call; will raise or be a no-op if unavailable)
        try:
            # Attempt to detect checkpointing state by trying to disable (no direct getter in HF API)
            model.gradient_checkpointing_disable()
            checkpoint_was_enabled = True
            # we will re-enable later
            print("[info] Temporarily disabled gradient checkpointing on model for .ft inference pass.")
        except Exception:
            checkpoint_was_enabled = False
    except Exception:
        checkpoint_was_enabled = False

    model.eval()
    inference_ctx = torch.inference_mode if hasattr(torch, "inference_mode") else torch.no_grad
    with inference_ctx():
        idx = 0
        for start in tqdm(range(0, n, batch_size), desc=f"Embedding (stream){suffix} {safe_name}"):
            batch_texts = lines[start:start+batch_size]
            emb_t = model_batch_embeddings_torch(batch_texts, tokenizer, model, max_length=BATCH_TOKEN_MAX_LEN)
            batch_embs = emb_t.detach().cpu().to(dtype=torch.float32).numpy()
            B = batch_embs.shape[0]

            # bulk write to memmap slice
            mm[idx:idx+B, :] = batch_embs

            # write text lines for this batch
            for b in range(B):
                row = batch_embs[b]
                txt_f.write(" ".join(f"{float(v):.6e}" for v in row) + "\n")

            idx += B

            # free intermediates and reduce fragmentation
            del emb_t
            del batch_embs
            if IS_CUDA:
                torch.cuda.empty_cache()

    # restore checkpointing if we disabled it
    if checkpoint_was_enabled:
        try:
            model.gradient_checkpointing_enable()
            print("[info] Re-enabled gradient checkpointing on model after .ft inference pass.")
        except Exception:
            pass

    txt_f.flush()
    txt_f.close()
    mm.flush()
    arr = np.load(emb_npy_path, mmap_mode='r')
    print(f"  Finished {device} (suffix={suffix}). .npy saved at {emb_npy_path} shape={arr.shape}")
    return arr

def load_embeddings_memmap_for_device(device, suffix=""):
    safe_name = device.replace('.json', '')
    emb_npy_path = os.path.join(EMB_DIR, f"{safe_name}.emb{suffix}.npy")
    total_lines = per_device_total_lines[device]
    if os.path.exists(emb_npy_path):
        try:
            arr = np.load(emb_npy_path, mmap_mode='r')
            if arr.ndim == 1:
                hidden_size = base_model.config.hidden_size
                if arr.size % hidden_size == 0:
                    rows = arr.size // hidden_size
                    arr = arr.reshape((rows, hidden_size))
            if arr.ndim == 2 and arr.shape[0] == total_lines:
                return arr
            else:
                return None
        except Exception:
            return None
    else:
        return None

# ---------------- Compute missing baseline embeddings up-front ----------------
def compute_missing_embeddings(devices):
    to_compute = []
    for device in devices:
        emb_arr = load_embeddings_memmap_for_device(device, suffix="")
        if emb_arr is None:
            to_compute.append(device)
        else:
            print(f"Skipping {device}: baseline .npy exists and looks valid ({emb_arr.shape[0]} rows).")
    for device in to_compute:
        run_with_oom_fallback(compute_and_cache_embeddings_stream, batch_arg_name="batch_size", start_batch=EMBED_BATCH_SIZE, min_batch=1, device=device)

# ---------------- helper: check finetuned embeddings exist ----------------
def all_finetuned_embeddings_exist(devices, suffix=".ft"):
    hidden_size = base_model.config.hidden_size
    for device in devices:
        emb_npy_path = os.path.join(EMB_DIR, f"{device.replace('.json','')}.emb{suffix}.npy")
        if not os.path.exists(emb_npy_path):
            return False
        try:
            arr = np.load(emb_npy_path, mmap_mode='r')
            if arr.ndim == 1:
                if arr.size % hidden_size == 0:
                    rows = arr.size // hidden_size
                    arr = arr.reshape((rows, hidden_size))
                else:
                    return False
            if arr.ndim != 2:
                return False
            if arr.shape[0] != per_device_total_lines[device]:
                return False
        except Exception:
            return False
    return True

def finetuned_model_files_exist(path):
    """Return True if this path looks like a HF-saved model directory (heuristic)."""
    if not os.path.isdir(path):
        return False
    candidates = [
        "pytorch_model.bin",
        "pytorch_model.safetensors",
        "tf_model.h5",
        "flax_model.msgpack",
        "adapter_model.bin",
    ]
    # also accept presence of config.json as a weak indicator
    if os.path.exists(os.path.join(path, "config.json")):
        # but make sure there is at least one weight file or safetensors
        for c in candidates:
            if os.path.exists(os.path.join(path, c)):
                return True
        # if only config.json exists, treat as present (you can make this stricter)
        return True
    return False

def load_finetuned_student_model(load_dir):
    """Load the finetuned model from load_dir into global student_model and put it on DEVICE."""
    global student_model
    print(f"[info] Loading finetuned student model from: {load_dir}", flush=True)
    if student_model is None:
        if IS_MAMBA:
            student_model = AutoModelForCausalLM.from_pretrained(load_dir, config=cfg, **dtype_kwargs)
        else:
            student_model = AutoModel.from_pretrained(load_dir, config=cfg, **dtype_kwargs)
        try:
            student_model.gradient_checkpointing_enable()
        except Exception:
            pass
    else:
        # already have an instance; reuse it
        pass
    student_model.to(DEVICE)
    student_model.eval()
    return student_model

# ---------------- plotting / confusion / pptx helpers ----------------
def normalize_cm_rows(cm):
    cm = cm.astype(float)
    row_sums = cm.sum(axis=1, keepdims=True)
    with np.errstate(divide='ignore', invalid='ignore'):
        normalized = np.divide(cm, row_sums, where=row_sums != 0)
    normalized[np.isnan(normalized)] = 0.0
    return normalized

def save_confusion_matrix_images(cm, class_names, out_base, title_text):
    fig, ax = plt.subplots(figsize=(8, 6))
    im = ax.imshow(cm, interpolation='nearest', vmin=0.0, vmax=1.0, cmap=plt.cm.Blues)
    ax.set_title(title_text)
    plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    n = len(class_names)
    ax.set_xticks(np.arange(n))
    ax.set_yticks(np.arange(n))
    ax.set_xticklabels(class_names, rotation=45, ha="right")
    ax.set_yticklabels(class_names)
    thresh = 0.5
    for i in range(cm.shape[0]):
        for j in range(cm.shape[1]):
            ax.text(j, i, f"{cm[i, j]:.2f}",
                    ha="center", va="center",
                    color="white" if cm[i, j] > thresh else "black", fontsize=8)
    plt.tight_layout()
    paths = {}
    for ext in ("png", "svg", "pdf"):
        p = f"{out_base}.{ext}"
        fig.savefig(p, dpi=300, transparent=True)
        paths[ext] = p
    plt.close(fig)
    return paths

def add_slide_with_image(prs, image_path, slide_title):
    if Presentation is None:
        raise RuntimeError("python-pptx is required to create PPTX. Install with `pip install python-pptx`")
    layout_index = 6 if len(prs.slide_layouts) > 6 else 5
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    left = Inches(0.5)
    top = Inches(0.2)
    width = Inches(9)
    height = Inches(0.6)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.font.size = Pt(18)
    img_left = Inches(0.5)
    img_top = Inches(1.0)
    img_width = Inches(9.0)
    try:
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
    except Exception as e:
        print(f"  Warning: could not insert image {image_path} into PPTX slide: {e}")

class FTData:
    def __init__(self, texts, labels, orig_embs, centroids):
        self.texts = texts
        self.labels = labels
        self.orig_embs = orig_embs.astype(np.float32)
        self.centroids = centroids
    def __len__(self):
        return len(self.texts)
    def __getitem__(self, idx):
        label = self.labels[idx]
        return self.texts[idx], int(label), self.orig_embs[idx], self.centroids[label]

# ---------------- Main execution ----------------
def main():
    global student_model   # use the single module-level student_model (avoid re-loading)

    # Compute baseline embeddings if missing
    compute_missing_embeddings(device_list)

    # ---------------- Build single finetune set (first FINE_PCT per device) ----------------
    per_device_fine_idxs = {}
    all_fine_texts = []
    all_fine_labels = []
    for device in device_list:
        total = per_device_total_lines[device]
        fine_count = int(total * FINE_PCT)
        fine_count = max(1, fine_count)
        fine_idxs = list(range(0, fine_count))
        per_device_fine_idxs[device] = fine_idxs
        lines = per_device_lines[device]
        labels = [device_label_map[device]] * len(fine_idxs)
        per_device_texts = [lines[i] for i in fine_idxs]
        all_fine_texts.append(per_device_texts)
        all_fine_labels.append(labels)

    per_device_texts_map = {device_list[i]: all_fine_texts[i] for i in range(len(device_list))}
    per_device_labels_map = {device_list[i]: all_fine_labels[i] for i in range(len(device_list))}
    fine_texts, fine_labels = interleave_round_robin(per_device_texts_map, device_list, device_label_map)

    print(f"\nPrepared global fine-tune set from first {FINE_PCT*100:.1f}% per device: total fine examples = {len(fine_texts)}")

    # ---------------- Train RF teacher on baseline embeddings of the FIRST fine block ----------------
    print("\nTraining fixed RF teacher on baseline embeddings from the FIRST fine block (assumption)...")
    X_teacher = []
    y_teacher = []
    for device in device_list:
        emb_arr = load_embeddings_memmap_for_device(device, suffix="")
        if emb_arr is None:
            emb_arr = run_with_oom_fallback(compute_and_cache_embeddings_stream, batch_arg_name="batch_size", start_batch=EMBED_BATCH_SIZE, min_batch=1, device=device)
        t_idxs = per_device_fine_idxs[device]
        if len(t_idxs) == 0:
            continue
        X_teacher.append(emb_arr[t_idxs, :])
        y_teacher.extend([device_label_map[device]] * len(t_idxs))
    if len(X_teacher) == 0:
        raise RuntimeError("No teacher training data found (unexpected).")
    X_teacher = np.vstack(X_teacher).astype(np.float32)
    y_teacher = np.array(y_teacher, dtype=int)
    print(f"  RF teacher training size: {X_teacher.shape[0]} examples")
    rf_teacher = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
    rf_teacher.fit(X_teacher, y_teacher)
    print("  RF teacher trained.")

    # ---------------- Decide whether to finetune or reuse existing finetuned embeddings (updated) ----------------
    skip_finetune = False
    if not FINETUNE_ENABLED:
        skip_finetune = True

    # If a saved finetuned model exists on disk we prefer to LOAD it and generate .ft embeddings
    finetuned_model_available = finetuned_model_files_exist(FINETUNED_MODEL_DIR)
    if FINETUNE_ENABLED and finetuned_model_available:
        # Load finetuned model (sets global student_model)
        try:
            load_finetuned_student_model(FINETUNED_MODEL_DIR)
            print("[info] Finetuned model found and loaded. We'll use it instead of running finetuning.")
            # Ensure .ft embeddings exist (or create them). honor REUSE_FINETUNE_EMBS_IF_PRESENT
            if REUSE_FINETUNE_EMBS_IF_PRESENT and all_finetuned_embeddings_exist(device_list, suffix=".ft"):
                print("[info] All .ft embeddings already exist and REUSE_FINETUNE_EMBS_IF_PRESENT=True -> skipping .ft recompute.")
            else:
                print("[info] Generating missing .ft embeddings using the loaded finetuned model...")
                # ensure student_model is on DEVICE before generating
                try:
                    student_model.to(DEVICE)
                except Exception:
                    pass
                for device in device_list:
                    # compute (with OOM fallback) only when missing (compute_and_cache checks and overwrites invalid files)
                    run_with_oom_fallback(
                        compute_and_cache_embeddings_stream_for_model,
                        batch_arg_name="batch_size",
                        start_batch=EMBED_BATCH_SIZE,
                        min_batch=1,
                        device=device,
                        model=student_model,
                        suffix=".ft"
                    )
            skip_finetune = True
        except Exception as e:
            print(f"[warning] Found finetuned model dir but failed to load it: {e}. Will attempt to finetune instead.", flush=True)
            skip_finetune = False
    else:
        # existing behavior: if .ft embeddings exist and user asked to reuse, skip finetune
        if FINETUNE_ENABLED and REUSE_FINETUNE_EMBS_IF_PRESENT and all_finetuned_embeddings_exist(device_list, suffix=".ft"):
            print("[info] All .ft embeddings exist and REUSE_FINETUNE_EMBS_IF_PRESENT=True -> skipping finetuning (will reuse embeddings).")
            skip_finetune = True

    # ---------------- STUDENT FINETUNE (single student trained once using fine_texts) ----------------
    if not skip_finetune:
        print("\n[1] Starting single student fine-tune using first-FINE_PCT interleaved examples (with RF feedback)...")

        print("[2] Loading teacher (baseline) embeddings for fine set from cache if available...")
        per_device_emb_slices = {}
        missing_cache = False

        for device in device_list:
            emb_arr = load_embeddings_memmap_for_device(device, suffix="")
            if emb_arr is None:
                missing_cache = True
                break
            idxs = per_device_fine_idxs[device]
            if len(idxs) == 0:
                per_device_emb_slices[device] = []
            else:
                per_device_emb_slices[device] = [np.array(emb_arr[i], dtype=np.float32) for i in idxs]

        if not missing_cache:
            orig_list, _ = interleave_round_robin(per_device_emb_slices, device_list, device_label_map)
            if len(orig_list) == 0:
                raise RuntimeError("No fine examples found when loading cached embeddings.")
            orig_cls_all = np.vstack(orig_list).astype(np.float32)
            print(f"[loaded] orig_cls_all shape: {orig_cls_all.shape}")
        else:
            print("[fallback] Some baseline .npy files missing/invalid -> computing teacher embeddings for fine_texts now.")
            orig_cls_parts = []
            for i in range(0, len(fine_texts), FT_BATCH_SIZE):
                batch = fine_texts[i:i+FT_BATCH_SIZE]
                emb_t = model_batch_embeddings_torch(batch, tokenizer, base_model, max_length=BATCH_TOKEN_MAX_LEN)
                orig_cls_parts.append(emb_t.cpu().numpy())
                del emb_t
                if IS_CUDA:
                    torch.cuda.empty_cache()
            orig_cls_all = np.vstack(orig_cls_parts).astype(np.float32)
            print(f"[computed] orig_cls_all shape: {orig_cls_all.shape}")
            try:
                del orig_cls_parts
            except Exception:
                pass

        print("[6] Creating labels array...")
        labels_arr = np.array(fine_labels, dtype=int)
        print(f"[6.1] labels_arr shape: {labels_arr.shape}", flush=True)

        print("[7] Computing centroids...")
        num_classes = len(device_list)
        centroids = np.zeros((num_classes, orig_cls_all.shape[1]), dtype=np.float32)
        for c in range(num_classes):
            idxs = np.where(labels_arr == c)[0]
            if len(idxs) > 0:
                centroids[c] = orig_cls_all[idxs].mean(axis=0)
        print(f"[7.1] Centroids computed for {num_classes} classes.", flush=True)

        print("[8] Creating fine-tune dataset and dataloader...")
        ft_dataset = FTData(fine_texts, fine_labels, orig_cls_all, centroids)
        ft_loader = DataLoader(
            ft_dataset,
            batch_size=FT_BATCH_SIZE,
            shuffle=False,
            pin_memory=IS_CUDA,
            num_workers= 0
        )
        print("[8.1] DataLoader ready.", flush=True)

        # Free base_model GPU memory before loading student model (important to avoid OOM)
        try:
            base_model.to("cpu")
            torch.cuda.empty_cache()
            print("[info] Moved base_model to CPU to free GPU for student.", flush=True)
        except Exception as e:
            print(f"[warning] Could not move base_model to CPU: {e}", flush=True)

        # Ensure single reusable student_model (module-level 'student_model' defined earlier)
        if student_model is None:
            print("[9] Loading student model (will create new instance)...", flush=True)
            if IS_MAMBA:
                student_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)
            else:
                student_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)

            try:
                student_model.gradient_checkpointing_enable()
                print("[info] Enabled gradient checkpointing on student_model.", flush=True)
            except Exception:
                pass

            print(f"[9.1] Model loaded. id(student_model)={id(student_model)}", flush=True)
        else:
            print(f"[9] Reusing existing student_model id={id(student_model)}", flush=True)

        student_model.to(DEVICE)
        student_model.train()

        print("[10.1] Model ready on device.", flush=True)

        print("[11] Creating classifier head and optimizer...")
        hidden_size = student_model.config.hidden_size
        classifier_head = nn.Linear(hidden_size, num_classes).to(DEVICE)

        mse_loss = nn.MSELoss()
        ce_loss_per_sample = nn.CrossEntropyLoss(reduction='none')
        optimizer = torch.optim.AdamW(list(student_model.parameters()) + list(classifier_head.parameters()), lr=LR)
        print("[11.1] Classifier and optimizer ready.", flush=True)

        print("  Fine-tune loop (anchor + proto + classification + RF-feedback)...", flush=True)

        for epoch in range(EPOCHS):
            print(f"[FT] Epoch {epoch+1}/{EPOCHS} start", flush=True)
            student_model.train()
            running_loss = 0.0
            steps = 0
            batch_idx = 0

            try:
                total_batches = len(ft_loader)
            except Exception:
                total_batches = None

            for batch in tqdm(ft_loader, desc=f" FT Epoch {epoch+1}/{EPOCHS}", leave=False):
                batch_idx += 1
                t_batch_start = time.time()
                print(f"[FT] epoch {epoch+1} batch {batch_idx}/{total_batches if total_batches is not None else '?'} - start", flush=True)

                try:
                    texts_b, labels_b, orig_emb_b, proto_emb_b = batch
                except Exception as e:
                    print(f"[FT][ERROR] failed to unpack batch: {e}", flush=True)
                    raise

                try:
                    n_texts = len(texts_b)
                except Exception:
                    n_texts = "?"
                print(f"[FT]  batch sizes -> texts: {n_texts}, labels: {len(labels_b)}", flush=True)

                try:
                    if isinstance(orig_emb_b, np.ndarray):
                        orig_for_rf = orig_emb_b
                    else:
                        orig_for_rf = orig_emb_b.cpu().numpy()
                    print(f"[FT]  orig_for_rf shape: {getattr(orig_for_rf,'shape', 'unknown')}", flush=True)
                except Exception as e:
                    print(f"[FT][ERROR] orig_for_rf conversion failed: {e}", flush=True)
                    raise

                try:
                    rf_preds = rf_teacher.predict(orig_for_rf)
                    rf_preds = np.array(rf_preds, dtype=int)
                    print(f"[FT]  rf_preds shape: {rf_preds.shape}, sample: {rf_preds[:5]}", flush=True)
                except Exception as e:
                    print(f"[FT][ERROR] rf_teacher.predict failed: {e}", flush=True)
                    raise

                labels_np = np.array(labels_b, dtype=int)

                # tokenization
                try:
                    t_tok_start = time.time()
                    enc = tokenizer(list(texts_b), return_tensors='pt', truncation=True, max_length=BATCH_TOKEN_MAX_LEN, padding=True)
                    enc = {k: v.to(DEVICE) for k, v in enc.items()}
                    t_tok_end = time.time()
                    shapes = {k: tuple(v.size()) for k, v in enc.items()}
                    print(f"[FT]  tokenized -> shapes: {shapes}, tokenization time: {(t_tok_end-t_tok_start):.3f}s", flush=True)
                except Exception as e:
                    print(f"[FT][ERROR] tokenization failed: {e}", flush=True)
                    raise

                # GPU memory snapshot (if CUDA)
                if IS_CUDA:
                    try:
                        alloc = torch.cuda.memory_allocated(DEVICE) / (1024**2)
                        reserved = torch.cuda.memory_reserved(DEVICE) / (1024**2) if hasattr(torch.cuda, "memory_reserved") else torch.cuda.memory_reserved() / (1024**2)
                        max_alloc = torch.cuda.max_memory_allocated(DEVICE) / (1024**2) if hasattr(torch.cuda, "max_memory_allocated") else torch.cuda.max_memory_allocated() / (1024**2)
                        print(f"[FT]  GPU mem (MiB): allocated={alloc:.1f}, reserved={reserved:.1f}, max_alloc={max_alloc:.1f}", flush=True)
                    except Exception as e:
                        print(f"[FT]  GPU mem query failed: {e}", flush=True)

                # forward (with grad) and pooling
                try:
                    t_fwd_start = time.time()
                    out = student_model(**enc)
                    t_fwd_end = time.time()
                    print(f"[FT]  forward done, time: {(t_fwd_end-t_fwd_start):.3f}s, out has_hidden_states={hasattr(out, 'hidden_states')}", flush=True)

                    if IS_MAMBA:
                        last_hidden = out.hidden_states[-1]  # (B, seq_len, hidden)
                        mask = enc["attention_mask"].unsqueeze(-1).to(last_hidden.dtype)  # (B, seq_len, 1)
                        summed = (last_hidden * mask).sum(dim=1)
                        counts = mask.sum(dim=1).clamp(min=1)
                        cls_emb = (summed / counts).to(dtype=torch.float32, device=DEVICE)
                    else:
                        cls_emb = out.last_hidden_state[:, 0, :].to(dtype=torch.float32, device=DEVICE)

                    print(f"[FT]  pooled cls_emb shape: {tuple(cls_emb.size())}", flush=True)
                except Exception as e:
                    print(f"[FT][ERROR] forward/pooling failed: {e}", flush=True)
                    raise

                # convert orig/proto to tensors
                try:
                    if isinstance(orig_emb_b, np.ndarray):
                        orig_emb_t = torch.from_numpy(orig_emb_b).to(dtype=torch.float32, device=DEVICE)
                    else:
                        orig_emb_t = orig_emb_b.to(dtype=torch.float32, device=DEVICE)

                    if isinstance(proto_emb_b, np.ndarray):
                        proto_emb_t = torch.from_numpy(proto_emb_b).to(dtype=torch.float32, device=DEVICE)
                    else:
                        proto_emb_t = proto_emb_b.to(dtype=torch.float32, device=DEVICE)

                    print(f"[FT]  orig_emb_t shape: {tuple(orig_emb_t.size())}, proto_emb_t shape: {tuple(proto_emb_t.size())}", flush=True)
                except Exception as e:
                    print(f"[FT][ERROR] orig/proto tensor conversion failed: {e}", flush=True)
                    raise

                # compute losses + backward
                try:
                    loss_anchor = mse_loss(cls_emb, orig_emb_t)
                    loss_proto = mse_loss(cls_emb, proto_emb_t)

                    logits = classifier_head(cls_emb)
                    labels_t = torch.tensor(labels_np, dtype=torch.long, device=DEVICE)
                    loss_cls_vec = ce_loss_per_sample(logits, labels_t)

                    rf_pred_t = torch.tensor(rf_preds, dtype=torch.long, device=DEVICE)
                    loss_rf_vec = ce_loss_per_sample(logits, rf_pred_t)
                    signs = torch.where(rf_pred_t == labels_t, 1.0, -1.0).to(DEVICE)
                    loss_rf = (loss_rf_vec * signs).mean()

                    loss = (LAMBDA_ANCHOR * loss_anchor +
                            LAMBDA_PROTO * loss_proto +
                            LAMBDA_CLS * loss_cls_vec.mean() +
                            LAMBDA_RF  * loss_rf)

                    print(f"[FT]  losses -> anchor: {float(loss_anchor.detach().cpu()):.6f}, proto: {float(loss_proto.detach().cpu()):.6f}, cls(mean): {float(loss_cls_vec.mean().detach().cpu()):.6f}, rf: {float(loss_rf.detach().cpu()):.6f}", flush=True)

                    if DEBUG_FT_VERBOSE and IS_CUDA:
                        try:
                            torch.cuda.synchronize(DEVICE)
                        except Exception:
                            pass

                    optimizer.zero_grad()
                    try:
                        t_bwd_start = time.time()
                        loss.backward()
                        if DEBUG_FT_VERBOSE and IS_CUDA:
                            torch.cuda.synchronize(DEVICE)
                        t_bwd_end = time.time()
                        print(f"[FT]  backward completed, time: {(t_bwd_end - t_bwd_start):.3f}s", flush=True)
                    except Exception as e:
                        print(f"[FT][ERROR] backward() raised: {e}", flush=True)
                        if IS_CUDA:
                            try:
                                print(f"[FT]  GPU mem snapshot BEFORE cleanup: allocated={torch.cuda.memory_allocated(DEVICE)/(1024**2):.1f} MiB reserved={torch.cuda.memory_reserved(DEVICE)/(1024**2):.1f} MiB", flush=True)
                            except Exception:
                                pass
                            torch.cuda.empty_cache()
                            print("[FT]  Emptied CUDA cache after backward error.", flush=True)
                        raise

                    try:
                        t_step_start = time.time()
                        optimizer.step()
                        if DEBUG_FT_VERBOSE and IS_CUDA:
                            torch.cuda.synchronize(DEVICE)
                        t_step_end = time.time()
                        print(f"[FT]  optimizer.step completed, time: {(t_step_end - t_step_start):.3f}s", flush=True)
                    except Exception as e:
                        print(f"[FT][ERROR] optimizer.step() raised: {e}", flush=True)
                        if IS_CUDA:
                            torch.cuda.empty_cache()
                        raise

                    running_loss += float(loss.detach().cpu().numpy())
                    steps += 1
                except Exception as e:
                    print(f"[FT][ERROR] loss/backward/step failed: {e}", flush=True)
                    raise

                t_batch_end = time.time()
                print(f"[FT] epoch {epoch+1} batch {batch_idx} done, elapsed {(t_batch_end-t_batch_start):.3f}s", flush=True)
                if IS_CUDA:
                    try:
                        alloc_post = torch.cuda.memory_allocated(DEVICE) / (1024**2)
                        reserved_post = torch.cuda.memory_reserved(DEVICE) / (1024**2)
                        print(f"[FT]  post-step GPU mem (MiB): allocated={alloc_post:.1f}, reserved={reserved_post:.1f}", flush=True)
                    except Exception:
                        pass

            avg_loss = running_loss / max(1, steps)
            la = float(loss_anchor.detach().cpu().numpy()) if isinstance(loss_anchor, torch.Tensor) else float(loss_anchor)
            lp = float(loss_proto.detach().cpu().numpy()) if isinstance(loss_proto, torch.Tensor) else float(loss_proto)
            lc = float(loss_cls_vec.mean().detach().cpu().numpy()) if isinstance(loss_cls_vec, torch.Tensor) else float(loss_cls_vec.mean())
            lrf = float(loss_rf.detach().cpu().numpy()) if isinstance(loss_rf, torch.Tensor) else float(loss_rf)
            print(f"   Epoch {epoch+1}/{EPOCHS} - avg loss: {avg_loss:.6f} (anchor {la:.6f}, proto {lp:.6f}, cls {lc:.6f}, rf {lrf:.6f})", flush=True)

        print("  Finished finetuning student.")

        # --- SAVE the finetuned model + tokenizer immediately (before long .ft embedding pass) ---
        try:
            os.makedirs(FINETUNED_MODEL_DIR, exist_ok=True)
            # move to CPU to avoid any CUDA-related state issues when saving
            try:
                student_model.to("cpu")
            except Exception:
                pass
            student_model.save_pretrained(FINETUNED_MODEL_DIR)
            tokenizer.save_pretrained(FINETUNED_MODEL_DIR)
            print(f"[info] Saved finetuned model + tokenizer to: {FINETUNED_MODEL_DIR}", flush=True)

            subject = "Model Finetuned!"
            body = "Model had been successfully finetuned"
            special.send_test_email(subject, body)

            # move model back to device for embedding generation
            try:
                student_model.to(DEVICE)
            except Exception:
                pass
        except Exception as e:
            print(f"[warning] Failed to save finetuned model to {FINETUNED_MODEL_DIR}: {e}", flush=True)

        # ---------------- Diagnostic comparison block ----------------
        # We'll compare embeddings from the in-memory student_model vs a reloaded copy
        # under a few configurations to surface differences.
        try:
            print("\n[DIAGNOSTIC] Starting embedding parity checks (in-memory vs reloaded saved model)...")
            # helper to prepare model for inference
            def prepare_model_for_inference(m, disable_checkpoint=True, to_device=None):
                if to_device is not None:
                    try:
                        m.to(to_device)
                    except Exception:
                        pass
                m.eval()
                if disable_checkpoint:
                    try:
                        m.gradient_checkpointing_disable()
                    except Exception:
                        pass

            def compare_models_on_texts(texts, tokenizer, model_a, model_b, device_a=None, device_b=None, label_a="A", label_b="B"):
                # compute embeddings for both models using specified devices (None -> use global DEVICE)
                da = device_a if device_a is not None else DEVICE
                db = device_b if device_b is not None else DEVICE
                # ensure models on their devices
                try:
                    model_a.to(da)
                except Exception:
                    pass
                try:
                    model_b.to(db)
                except Exception:
                    pass
                # prepare models (disable checkpointing & eval)
                try:
                    model_a.eval()
                    model_a.gradient_checkpointing_disable()
                except Exception:
                    pass
                try:
                    model_b.eval()
                    model_b.gradient_checkpointing_disable()
                except Exception:
                    pass

                # use inference mode for both if available
                inference_ctx = torch.inference_mode if hasattr(torch, "inference_mode") else torch.no_grad
                with inference_ctx():
                    if da.type == "cpu":
                        emb_a = get_embedding_batch_on_device(texts, tokenizer, model_a, device=da, max_length=BATCH_TOKEN_MAX_LEN)
                    else:
                        emb_a = get_embedding_batch(texts, tokenizer, model_a, max_length=BATCH_TOKEN_MAX_LEN)
                    if db.type == "cpu":
                        emb_b = get_embedding_batch_on_device(texts, tokenizer, model_b, device=db, max_length=BATCH_TOKEN_MAX_LEN)
                    else:
                        emb_b = get_embedding_batch(texts, tokenizer, model_b, max_length=BATCH_TOKEN_MAX_LEN)

                diffs = emb_a - emb_b
                max_abs = float(np.max(np.abs(diffs)))
                mean_abs = float(np.mean(np.abs(diffs)))
                rmse = float(np.sqrt(np.mean(diffs**2)))
                denom = (np.max(np.abs(emb_a)) + 1e-12)
                max_rel = max_abs / denom
                print(f"[DIAG] {label_a} vs {label_b} | device {da} / {db} -> max_abs={max_abs:.6e}, mean_abs={mean_abs:.6e}, rmse={rmse:.6e}, max_rel={max_rel:.6e}")
                print(f"[DIAG] sample diffs first vector (first 8): {diffs[0][:8].tolist()}")
                return emb_a, emb_b, diffs

            # pick a small list of sample texts (few examples from first device)
            sample_texts = []
            first_dev = device_list[0]
            sample_count = min(5, len(per_device_lines[first_dev]))
            for i in range(sample_count):
                sample_texts.append(per_device_lines[first_dev][i])

            # reload model freshly from disk
            try:
                if IS_MAMBA:
                    reloaded = AutoModelForCausalLM.from_pretrained(FINETUNED_MODEL_DIR, config=cfg, **dtype_kwargs)
                else:
                    reloaded = AutoModel.from_pretrained(FINETUNED_MODEL_DIR, config=cfg, **dtype_kwargs)
                reloaded.to(DEVICE)
                reloaded.eval()
                try:
                    reloaded.gradient_checkpointing_disable()
                except Exception:
                    pass
            except Exception as e:
                print(f"[DIAG] Failed to reload saved model for diagnostics: {e}")
                reloaded = None

            # Test 1: both on same DEVICE, both with checkpointing disabled
            if reloaded is not None:
                print("[DIAG] Test 1: both on same DEVICE, checkpointing disabled (recommended)")
                prepare_model_for_inference(student_model, disable_checkpoint=True, to_device=DEVICE)
                prepare_model_for_inference(reloaded, disable_checkpoint=True, to_device=DEVICE)
                compare_models_on_texts(sample_texts, tokenizer, student_model, reloaded, device_a=DEVICE, device_b=DEVICE, label_a="in-memory", label_b="reloaded")

                # Test 2: in-memory keep checkpointing enabled (simulate training state) vs reloaded disabled
                print("[DIAG] Test 2: in-memory checkpointing ENABLED (train-like) vs reloaded checkpointing DISABLED")
                try:
                    student_model.gradient_checkpointing_enable()
                except Exception:
                    pass
                # set eval but keep checkpointing hooks
                student_model.eval()
                try:
                    reloaded.gradient_checkpointing_disable()
                except Exception:
                    pass
                compare_models_on_texts(sample_texts, tokenizer, student_model, reloaded, device_a=DEVICE, device_b=DEVICE, label_a="in-memory(checkpoint_on)", label_b="reloaded")

                # Test 3: both on CPU (move to CPU for maximal determinism)
                print("[DIAG] Test 3: move both to CPU and compare")
                prepare_model_for_inference(student_model, disable_checkpoint=True, to_device=torch.device("cpu"))
                prepare_model_for_inference(reloaded, disable_checkpoint=True, to_device=torch.device("cpu"))
                compare_models_on_texts(sample_texts, tokenizer, student_model, reloaded, device_a=torch.device("cpu"), device_b=torch.device("cpu"), label_a="in-memory(CPU)", label_b="reloaded(CPU)")

                # After CPU test, move reloaded back to DEVICE (student_model will be moved back later)
                try:
                    reloaded.to(DEVICE)
                    reloaded.eval()
                    try:
                        reloaded.gradient_checkpointing_disable()
                    except Exception:
                        pass
                except Exception:
                    pass

                # Restore student_model to DEVICE and disable checkpointing for main flow
                try:
                    student_model.to(DEVICE)
                    student_model.eval()
                    try:
                        student_model.gradient_checkpointing_disable()
                    except Exception:
                        pass
                except Exception:
                    pass

            else:
                print("[DIAG] Skipping diagnostics because reloaded model could not be created.")
            print("[DIAGNOSTIC] Done.\n")
        except Exception as e:
            print(f"[DIAGNOSTIC] Diagnostic run failed: {e}", flush=True)

        print("  Saving finetuned embeddings (.ft) per device...")

        # Important: ensure we try a reasonably large start_batch and OOM-resize if needed
        for device in device_list:
            run_with_oom_fallback(
                compute_and_cache_embeddings_stream_for_model,
                batch_arg_name="batch_size",
                start_batch=EMBED_BATCH_SIZE,
                min_batch=1,
                device=device,
                model=student_model,
                suffix=".ft"
            )

        print("  Finetuned embeddings saved.")
    else:
        print("\nSkipping finetuning - will use existing .ft embeddings (if present) or skip finetuned evaluation entirely.")

    # ---------------- RF sweep training/eval using baseline and finetuned embeddings ----------------
    results_baseline = []
    results_finetuned = []
    cm_baseline_records = []
    cm_finetuned_records = []
    classes = [d.replace('.json', '') for d in device_list]

    rf_train_pcts = np.arange(START_PCT, END_PCT + 1e-9, STEP_PCT)
    for pct in rf_train_pcts:
        RF_TRAIN_PCT = float(pct)
        print("------------------------------------------------------------")
        print(f"Running sweep step: RF_TRAIN_PCT = {RF_TRAIN_PCT*100:.1f}% (RF-val fixed at last {RF_VAL_PCT*100:.1f}%)")
        per_device_train_emb_idxs = {}
        per_device_val_emb_idxs = {}
        total_train = 0
        total_val = 0
        for device in device_list:
            total_lines = per_device_total_lines[device]
            rf_val_count = int(total_lines * RF_VAL_PCT)
            rf_val_count = max(1, rf_val_count)

            fine_end_idx = int(total_lines * FINE_PCT)
            fine_end_idx = min(fine_end_idx, total_lines - rf_val_count)
            if fine_end_idx < 0:
                fine_end_idx = 0

            fine_end_idx = int(total_lines * RESERVED_FINE_PCT)
            fine_end_idx = min(fine_end_idx, total_lines - rf_val_count)
            if fine_end_idx < 0:
                fine_end_idx = 0

            rf_train_start_idx = fine_end_idx
            # compute train end as proportion of the whole file, then clamp to avoid overlapping the val tail
            rf_train_end_idx = int(total_lines * RF_TRAIN_PCT)
            rf_train_end_idx = min(rf_train_end_idx, total_lines - rf_val_count)
            if rf_train_end_idx <= rf_train_start_idx:
                train_idxs = []
            else:
                train_idxs = list(range(rf_train_start_idx, rf_train_end_idx))


            val_idxs = list(range(total_lines - rf_val_count, total_lines))

            per_device_train_emb_idxs[device] = train_idxs
            per_device_val_emb_idxs[device] = val_idxs
            total_train += len(train_idxs)
            total_val += len(val_idxs)

            if PRINT_SAMPLE_LINES:
                lines = per_device_lines[device]
                t_first = short(lines[train_idxs[0]]) if train_idxs else "<none>"
                t_last = short(lines[train_idxs[-1]]) if train_idxs else "<none>"
                v_first = short(lines[val_idxs[0]]) if val_idxs else "<none>"
                v_last = short(lines[val_idxs[-1]]) if val_idxs else "<none>"
                fine_first = short(lines[0]) if fine_end_idx > 0 else "<none>"
                fine_last = short(lines[fine_end_idx-1]) if fine_end_idx > 0 else "<none>"
                print(f"  {device}: total={total_lines}, fine 1-based=(1,{fine_end_idx}) count={fine_end_idx}, train 1-based=({rf_train_start_idx+1},{rf_train_end_idx}) count={len(train_idxs)}, val 1-based=({total_lines-len(val_idxs)+1},{total_lines}) count={len(val_idxs)}")
                if fine_end_idx > 0:
                    print(f"    fine sample first/last: {fine_first} / {fine_last}")
                if train_idxs:
                    print(f"    train sample first/last: {t_first} / {t_last}")
                if val_idxs:
                    print(f"    val   sample first/last: {v_first} / {v_last}")

        print(f"  Combined totals before interleave: n_train={total_train}, n_val={total_val}")

        # Baseline
        per_device_train_lists = {}
        per_device_val_lists = {}
        for device in device_list:
            emb_arr = load_embeddings_memmap_for_device(device, suffix="")
            if emb_arr is None:
                emb_arr = run_with_oom_fallback(compute_and_cache_embeddings_stream, batch_arg_name="batch_size", start_batch=EMBED_BATCH_SIZE, min_batch=1, device=device)
            t_idxs = per_device_train_emb_idxs[device]
            v_idxs = per_device_val_emb_idxs[device]
            per_device_train_lists[device] = [np.array(emb_arr[i], dtype=np.float32) for i in t_idxs] if len(t_idxs) > 0 else []
            per_device_val_lists[device]   = [np.array(emb_arr[i], dtype=np.float32) for i in v_idxs] if len(v_idxs) > 0 else []

        train_emb_list, train_labels = interleave_round_robin(per_device_train_lists, device_list, device_label_map)
        val_emb_list, val_labels     = interleave_round_robin(per_device_val_lists, device_list, device_label_map)

        if len(train_emb_list) == 0:
            print("  No baseline training samples available for this pct; skipping baseline.")
        else:
            X_train = np.vstack(train_emb_list).astype(np.float32)
            y_train = np.array(train_labels, dtype=int)
            X_val = np.vstack(val_emb_list).astype(np.float32) if len(val_emb_list) > 0 else np.zeros((0, X_train.shape[1]), dtype=np.float32)
            y_val = np.array(val_labels, dtype=int) if len(val_emb_list) > 0 else np.array([], dtype=int)

            print(f"  Baseline combined totals after interleave: n_train={X_train.shape[0]}, n_val={X_val.shape[0]}")
            rf = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
            rf.fit(X_train, y_train)
            if X_val.shape[0] == 0:
                print("  No baseline validation samples; skipping baseline eval.")
            else:
                y_pred = rf.predict(X_val)
                f1m = f1_score(y_val, y_pred, average='macro')
                print(f"  -> Baseline Macro F1 at RF_TRAIN_PCT={RF_TRAIN_PCT*100:.1f}%: {f1m:.4f}")
                cm = confusion_matrix(y_val, y_pred, labels=list(range(len(device_list))))
                cm_norm = normalize_cm_rows(cm)
                out_base = os.path.join(CM_BASELINE_DIR, f"confusion_baseline_RFtrain_{int(RF_TRAIN_PCT*100):02d}pct")
                title_text = f"Baseline RF Train {RF_TRAIN_PCT*100:.1f}% — Macro F1: {f1m:.4f}"
                saved = save_confusion_matrix_images(cm_norm, classes, out_base, title_text)
                cm_baseline_records.append({"pct": RF_TRAIN_PCT, "macro_f1": float(f1m), "n_train": int(X_train.shape[0]), "n_val": int(X_val.shape[0]), "img_png": saved.get("png"), "cm": cm_norm})
                results_baseline.append({"rf_train_pct": RF_TRAIN_PCT, "macro_f1": float(f1m), "n_train": int(X_train.shape[0]), "n_val": int(X_val.shape[0])})

        # Finetuned
        if FINETUNE_ENABLED:
            per_device_train_lists_ft = {}
            per_device_val_lists_ft = {}
            for device in device_list:
                emb_arr_ft = load_embeddings_memmap_for_device(device, suffix=".ft")
                if emb_arr_ft is None:
                    if student_model is None:
                        if IS_MAMBA:
                            student_model = AutoModelForCausalLM.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)
                        else:
                            student_model = AutoModel.from_pretrained(MODEL_NAME, config=cfg, **dtype_kwargs)
                        student_model.to(DEVICE)
                    emb_arr_ft = run_with_oom_fallback(compute_and_cache_embeddings_stream_for_model, batch_arg_name="batch_size", start_batch=EMBED_BATCH_SIZE, min_batch=1, device=device, model=student_model, suffix=".ft")
                t_idxs = per_device_train_emb_idxs[device]
                v_idxs = per_device_val_emb_idxs[device]
                per_device_train_lists_ft[device] = [np.array(emb_arr_ft[i], dtype=np.float32) for i in t_idxs] if len(t_idxs) > 0 else []
                per_device_val_lists_ft[device]   = [np.array(emb_arr_ft[i], dtype=np.float32) for i in v_idxs] if len(v_idxs) > 0 else []

            train_emb_list_ft, train_labels_ft = interleave_round_robin(per_device_train_lists_ft, device_list, device_label_map)
            val_emb_list_ft, val_labels_ft     = interleave_round_robin(per_device_val_lists_ft, device_list, device_label_map)

            if len(train_emb_list_ft) == 0:
                print("  No finetuned training samples available for this pct; skipping finetuned RF.")
            else:
                X_train_ft = np.vstack(train_emb_list_ft).astype(np.float32)
                y_train_ft = np.array(train_labels_ft, dtype=int)
                X_val_ft = np.vstack(val_emb_list_ft).astype(np.float32) if len(val_emb_list_ft) > 0 else np.zeros((0, X_train_ft.shape[1]), dtype=np.float32)
                y_val_ft = np.array(val_labels_ft, dtype=int) if len(val_emb_list_ft) > 0 else np.array([], dtype=int)

                print(f"  Finetuned combined totals after interleave: n_train={X_train_ft.shape[0]}, n_val={X_val_ft.shape[0]}")
                rf_ft = RandomForestClassifier(n_estimators=RF_N_ESTIMATORS, n_jobs=RF_N_JOBS, random_state=42)
                rf_ft.fit(X_train_ft, y_train_ft)
                if X_val_ft.shape[0] == 0:
                    print("  No finetuned validation samples; skipping finetuned eval.")
                else:
                    y_pred_ft = rf_ft.predict(X_val_ft)
                    f1m_ft = f1_score(y_val_ft, y_pred_ft, average='macro')
                    print(f"  -> Finetuned Macro F1 at RF_TRAIN_PCT={RF_TRAIN_PCT*100:.1f}%: {f1m_ft:.4f}")
                    cm_ft = confusion_matrix(y_val_ft, y_pred_ft, labels=list(range(len(device_list))))
                    cm_ft_norm = normalize_cm_rows(cm_ft)
                    out_base_ft = os.path.join(CM_FINETUNED_DIR, f"confusion_finetuned_RFtrain_{int(RF_TRAIN_PCT*100):02d}pct")
                    title_text_ft = f"Finetuned RF Train {RF_TRAIN_PCT*100:.1f}% — Macro F1: {f1m_ft:.4f}"
                    saved_ft = save_confusion_matrix_images(cm_ft_norm, classes, out_base_ft, title_text_ft)
                    cm_finetuned_records.append({"pct": RF_TRAIN_PCT, "macro_f1": float(f1m_ft), "n_train": int(X_train_ft.shape[0]), "n_val": int(X_val_ft.shape[0]), "img_png": saved_ft.get("png"), "cm": cm_ft_norm})
                    results_finetuned.append({"rf_train_pct": RF_TRAIN_PCT, "macro_f1": float(f1m_ft), "n_train": int(X_train_ft.shape[0]), "n_val": int(X_val_ft.shape[0])})
        else:
            print("  Finetuned evaluation skipped (FINETUNE_ENABLED=False).")

    # Save numeric results and create PPTX (same as before)
    def save_results_csv_and_plot(results_list, csv_name, plot_name):
        if len(results_list) == 0:
            return None
        csv_path = csv_name
        with open(csv_path, "w", newline='', encoding='utf-8') as csvf:
            writer = csv.DictWriter(csvf, fieldnames=["rf_train_pct", "n_train", "n_val", "macro_f1"])
            writer.writeheader()
            for r in results_list:
                writer.writerow({"rf_train_pct": r["rf_train_pct"], "n_train": r["n_train"], "n_val": r["n_val"], "macro_f1": r["macro_f1"]})
        pcts = [r["rf_train_pct"] * 100.0 for r in results_list]
        f1s = [r["macro_f1"] for r in results_list]
        plt.figure(figsize=(8, 5))
        plt.plot(pcts, f1s, marker='o')
        plt.xlabel("RF_TRAIN_PCT (%)")
        plt.ylabel("Macro F1 (validation)")
        plt.title(plot_name)
        plt.grid(True)
        plt.tight_layout()
        plot_path = csv_name.replace(".csv", ".png")
        plt.savefig(plot_path, dpi=300, transparent=True)
        plt.close()
        return csv_path, plot_path

    base_csv, base_plot = None, None
    if results_baseline:
        base_csv, base_plot = save_results_csv_and_plot(results_baseline, "rf_train_pct_results_baseline.csv", "Baseline: RF training size vs Macro F1")
        print(f"\nSaved baseline numeric results to: {base_csv}")

    ft_csv, ft_plot = None, None
    if FINETUNE_ENABLED and results_finetuned:
        ft_csv, ft_plot = save_results_csv_and_plot(results_finetuned, "rf_train_pct_results_finetuned.csv", "Finetuned: RF training size vs Macro F1")
        print(f"\nSaved finetuned numeric results to: {ft_csv}")
    elif not FINETUNE_ENABLED:
        print("\nFinetuning disabled -> no finetuned results saved.")

    def create_pptx_from_records(records, out_pptx):
        if Presentation is None:
            print(f"Skipping PPTX creation ({out_pptx}): python-pptx not installed.")
            return
        prs = Presentation()
        try:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = "Confusion Matrices"
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated: {datetime.utcnow().isoformat()} UTC"
        except Exception:
            pass
        for rec in records:
            title = f"RF Train {rec['pct']*100:.1f}% — Macro F1: {rec['macro_f1']:.4f} — n_train={rec['n_train']} n_val={rec['n_val']}"
            img = rec.get("img_png")
            if img is None:
                layout_index = 1 if len(prs.slide_layouts) > 1 else 0
                s = prs.slides.add_slide(prs.slide_layouts[layout_index])
                tb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1))
                tb.text_frame.text = title
                continue
            add_slide_with_image(prs, img, title)
        prs.save(out_pptx)
        print(f"Saved PPTX: {out_pptx}")

    pptx_base_path = None
    pptx_ft_path = None
    if cm_baseline_records:
        pptx_base_path = os.path.join(CM_DIR, f"confusion_baseline_{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}.pptx")
        create_pptx_from_records(cm_baseline_records, pptx_base_path)
    if FINETUNE_ENABLED and cm_finetuned_records:
        pptx_ft_path = os.path.join(CM_DIR, f"confusion_finetuned_{datetime.utcnow().strftime('%Y%m%dT%H%M%SZ')}.pptx")
        create_pptx_from_records(cm_finetuned_records, pptx_ft_path)

    print("\nDone.")
    summary = {
        "baseline_csv": base_csv,
        "baseline_plot": base_plot,
        "finetuned_csv": ft_csv if FINETUNE_ENABLED else None,
        "finetuned_plot": ft_plot if FINETUNE_ENABLED else None,
        "pptx_baseline": pptx_base_path,
        "pptx_finetuned": pptx_ft_path if FINETUNE_ENABLED else None,
        "n_baseline_results": len(results_baseline),
        "n_finetuned_results": len(results_finetuned) if FINETUNE_ENABLED else 0,
        "finetune_enabled": FINETUNE_ENABLED,
    }
    return summary

# ---------------- Run main inside wrapper and send emails on success/failure ----------------
if __name__ == "__main__":
    try:
        summary = main()
        subject = "Code Completed Successfully"
        body_lines = [
            "Your RF sweep + finetune script completed successfully.",
            "",
            f"Model: {MODEL_NAME}",
            f"Finetune enabled: {summary.get('finetune_enabled', False)}",
            f"Baseline results points: {summary.get('n_baseline_results', 0)}",
            f"Finetuned results points: {summary.get('n_finetuned_results', 0)}",
            ""
        ]
        if summary.get("baseline_csv"):
            body_lines.append(f"Baseline CSV: {summary['baseline_csv']}")
        if summary.get("baseline_plot"):
            body_lines.append(f"Baseline plot: {summary['baseline_plot']}")
        if summary.get("finetuned_csv"):
            body_lines.append(f"Finetuned CSV: {summary['finetuned_csv']}")
        if summary.get("finetuned_plot"):
            body_lines.append(f"Finetuned plot: {summary['finetuned_plot']}")
        if summary.get("pptx_baseline"):
            body_lines.append(f"Baseline PPTX: {summary['pptx_baseline']}")
        if summary.get("pptx_finetuned"):
            body_lines.append(f"Finetuned PPTX: {summary['pptx_finetuned']}")
        body = "\n".join(body_lines)
        try:
            special.send_test_email(subject, body)
        except Exception as send_e:
            print(f"Warning: sending success email failed: {send_e}")
            print("Success email body would have been:\n", body)
    except Exception as e:
        error_subject = "Code Failed!"
        tb_str = traceback.format_exc()
        error_body = f"The code encountered an error:\n{str(e)}\n\nFull traceback:\n{tb_str}"
        try:
            special.send_test_email(error_subject, error_body)
        except Exception as send_e:
            print(f"Warning: sending error email failed: {send_e}")
            print("Error body would have been:\n", error_body)
        raise
